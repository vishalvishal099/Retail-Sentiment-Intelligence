"""
Generate the Post-Midsem Project Presentation for Retail Sentiment Intelligence.
Creates a professional PPTX with ~35 slides covering:
  1. Title & Agenda
  2. System Architecture & Pipeline
  3. Vision/Image Processing Challenge & Mitigation
  4. Trust Score & Confidence Calculations
  5. ModernBERT Training Journey
  6. Dashboard & Data Population
  7. Notification System
  8. Review & Validate
  9. Post Explorer & Post Lifecycle
  10. Insights & Conclusion
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Configuration ───────────────────────────────────────────────────────────
OUTPUT_PATH = os.path.join(os.path.dirname(__file__), "..", "docs", "Sem_4",
                           "RSI_Post_Midsem_Presentation.pptx")
OUTPUT_PATH = os.path.normpath(OUTPUT_PATH)

# Brand colors
WALMART_BLUE = RGBColor(0x00, 0x71, 0xDC)
DARK_BLUE = RGBColor(0x04, 0x1E, 0x42)
LIGHT_BLUE = RGBColor(0xE8, 0xF4, 0xFD)
ACCENT_YELLOW = RGBColor(0xFF, 0xC2, 0x20)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
GREEN = RGBColor(0x10, 0xB9, 0x81)
RED = RGBColor(0xEF, 0x44, 0x44)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_box(slide, text, top=Inches(0.4), left=Inches(0.5),
                  width=Inches(9), height=Inches(0.8),
                  font_size=Pt(28), bold=True, color=DARK_BLUE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    return txBox


def add_body_box(slide, text, top=Inches(1.4), left=Inches(0.5),
                 width=Inches(9), height=Inches(5.5),
                 font_size=Pt(14), color=DARK_GRAY, bold=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = font_size
        p.font.color.rgb = color
        p.font.bold = bold
        p.space_after = Pt(6)
    return txBox


def add_bullet_slide(slide, title, bullets, sub_bullets=None):
    add_title_box(slide, title)
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)
        p.level = 0

        if sub_bullets and i in sub_bullets:
            for sb in sub_bullets[i]:
                p2 = tf.add_paragraph()
                p2.text = sb
                p2.font.size = Pt(13)
                p2.font.color.rgb = MED_GRAY
                p2.space_after = Pt(4)
                p2.level = 1


def add_table_slide(slide, title, headers, rows, top=Inches(1.5)):
    add_title_box(slide, title)
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_width = Inches(9)
    tbl_height = Inches(0.4) * n_rows
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(0.5), top, tbl_width, tbl_height)
    table = table_shape.table

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(12)
            para.font.bold = True
            para.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = WALMART_BLUE

    # Data rows
    for i, row in enumerate(rows, 1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(val)
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(11)
                para.font.color.rgb = DARK_GRAY
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY


def add_section_divider(prs, section_num, section_title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, DARK_BLUE)
    # Section number
    add_title_box(slide, f"0{section_num}" if section_num < 10 else str(section_num),
                  top=Inches(2.0), left=Inches(0.8), font_size=Pt(60),
                  color=ACCENT_YELLOW, bold=True)
    # Section title
    add_title_box(slide, section_title,
                  top=Inches(3.2), left=Inches(0.8), font_size=Pt(32),
                  color=WHITE, bold=True)
    if subtitle:
        add_body_box(slide, subtitle, top=Inches(4.2), left=Inches(0.8),
                     font_size=Pt(16), color=RGBColor(0xBB, 0xBB, 0xBB))
    return slide


# ═══════════════════════════════════════════════════════════════════════════════
#                          BUILD PRESENTATION
# ═══════════════════════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ─── SLIDE 1: Title ──────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BLUE)
add_title_box(slide, "Retail Sentiment Intelligence",
              top=Inches(2.0), left=Inches(0.8), width=Inches(8.5),
              font_size=Pt(36), color=WHITE)
add_body_box(slide,
             "Real-Time Brand Health Monitoring via Reddit NLP Pipeline\n\n"
             "Post Mid-Semester Presentation — Semester 4\n\n"
             "Vishal Singh\n"
             "M.Tech CSE — Lovely Professional University",
             top=Inches(3.2), left=Inches(0.8), font_size=Pt(18), color=RGBColor(0xCC, 0xCC, 0xCC))

# ─── SLIDE 2: Agenda ─────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_box(slide, "Agenda")
bullets = [
    "1.  System Architecture & End-to-End Pipeline",
    "2.  Vision/Image Processing — Challenge & Mitigation",
    "3.  Trust Score & Confidence Calculations",
    "4.  ModernBERT — Domain Fine-Tuning Journey",
    "5.  Dashboard — Data Population & Sections",
    "6.  Notification System — Group-Based Routing",
    "7.  Review & Validate — Human-in-the-Loop",
    "8.  Post Explorer & Filtering",
    "9.  Post Lifecycle — Kanban Workflow",
    "10. Insights & Competitor Analysis",
    "11. Results & Future Work",
]
add_bullet_slide(slide, "Agenda", bullets)

# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 1: ARCHITECTURE & PIPELINE
# ═══════════════════════════════════════════════════════════════════════════════
add_section_divider(prs, 1, "System Architecture & Pipeline",
                    "From Reddit data to actionable brand intelligence")

# Slide: High-Level Architecture
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_box(slide, "High-Level System Architecture")
add_body_box(slide,
    "┌─────────────────────────────────────────────────────────────────┐\n"
    "│  DATA SOURCES   →   PIPELINE (6 stages)   →   DASHBOARD (React) │\n"
    "└─────────────────────────────────────────────────────────────────┘\n\n"
    "Data Sources:  Arctic Shift API (free, no auth) + PRAW (optional)\n"
    "Pipeline:      Python — asyncio scheduler, 60-min + on-demand\n"
    "Storage:       SQLite (local dev) / Azure Cosmos DB (prod)\n"
    "API:           FastAPI + WebSocket — port 8001\n"
    "Frontend:      React 18 + TypeScript + Vite + Tailwind — port 5173\n"
    "LLM Runtime:   HuggingFace (offline) + Ollama (vision) + Azure OpenAI (opt-in)",
    top=Inches(1.4), font_size=Pt(14))

# Slide: 6-Stage Pipeline
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_box(slide, "6-Stage Pipeline Architecture")
add_table_slide(slide, "6-Stage Pipeline Architecture",
    ["Stage", "Component", "Purpose", "Key Technology"],
    [
        ["1. INGEST", "arctic_shift.py", "Fetch Reddit posts (hourly)", "Arctic Shift API, PRAW"],
        ["2. PREPROCESS", "preprocess.py", "Clean, langdetect, dedup", "langdetect, MiniLM-L6-v2"],
        ["3. TRUST", "trust/scorer.py", "Metadata + dedup + LLM credibility", "Weighted formula (0.4/0.3/0.3)"],
        ["4. ANALYZE", "analyzer.py", "Sentiment + Aspect + Vision", "ModernBERT, BART-MNLI, Gemma3"],
        ["5. AGGREGATE", "aggregator.py", "Hourly & daily rollups", "SQL window aggregation"],
        ["6. ALERT", "alerts/engine.py", "Spike & severity detection", "σ-threshold + rules"],
    ])

# Slide: Pipeline Sequence
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_box(slide, "Pipeline Execution Flow")
add_body_box(slide,
    "Scheduler (60-min tick or manual 'Run Now')\n"
    "    │\n"
    "    ├──▶ INGEST: Arctic Shift → fetch posts after cursor\n"
    "    │       └── Incremental: only new posts (cursor-based pagination)\n"
    "    ├──▶ PREPROCESS: Clean HTML, detect English, semantic dedup\n"
    "    │       └── MiniLM-L6-v2 embeddings → cosine sim > 0.92 = duplicate\n"
    "    ├──▶ TRUST SCORE: Metadata heuristics + LLM credibility check\n"
    "    │       └── Low-trust posts FLAGGED (not dropped) per requirements\n"
    "    ├──▶ ANALYZE: Sentiment (ModernBERT) + Aspects (BART-MNLI)\n"
    "    │       └── Vision: image posts → Gemma3:4b multi-pass → caption\n"
    "    ├──▶ AGGREGATE: Hourly/daily rollups → brand health KPIs\n"
    "    └──▶ ALERT: Spike detection (>2σ) + severity rules → notify\n\n"
    "Total pipeline latency: ~3–5 min for 25 subreddits\n"
    "25 tracked subreddits across 6 groups (Walmart core, Spark, Pharmacy, etc.)",
    top=Inches(1.3), font_size=Pt(13))

# Slide: Data Model
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_box(slide, "Storage Schema — SQLite / Cosmos DB")
add_table_slide(slide, "Storage Schema — SQLite / Cosmos DB",
    ["Table", "Partition Key", "Key Fields", "Purpose"],
    [
        ["raw_posts", "/subreddit", "id, title, body, author_hash, created_utc", "Ingested data (privacy-safe)"],
        ["analyses", "/subreddit", "post_id, sentiment, confidence, aspects, trust_score", "AI analysis results"],
        ["aggregates", "/time_window", "subreddit, window, metrics_json", "Pre-computed KPIs"],
        ["feedback", "/analyst_id", "post_id, correction, reply_text", "Human corrections + replies"],
        ["alerts", "/severity", "type, aspect, threshold_breached", "Triggered anomalies"],
        ["notification_groups", "id", "subreddits[], email_dl[], priority_filter", "Notification routing config"],
        ["notification_log", "group_id", "post_id, channel, status, sent_at", "Delivery audit trail"],
    ])

# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 2: VISION / IMAGE PROCESSING
# ═══════════════════════════════════════════════════════════════════════════════
add_section_divider(prs, 2, "Vision / Image Processing",
                    "Eliminating hallucination in multimodal retail analysis")

# Slide: The Problem
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_box(slide, "The Problem: Image-Only Posts")
add_body_box(slide,
    "• 3.9% of Reddit posts have empty bodies — complaint lives in the IMAGE\n"
    "• Screenshots of errors, damaged products, receipts, app glitches\n"
    "• Text-only pipeline misses these entirely → incomplete brand health\n\n"
    "Example:\n"
    "  Title: \"Can anyone help me? I need this fixed\"\n"
    "  Body:  (empty)\n"
    "  Image: Screenshot of Walmart app error message\n\n"
    "Without vision AI → this post is scored on title alone (useless)\n"
    "With vision AI → we extract the error text and understand the complaint",
    top=Inches(1.3), font_size=Pt(15))

# Slide: Model Selection
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_table_slide(slide, "Vision Model Selection — Why Gemma 3 4B",
    ["Model", "DocVQA", "Size", "Ollama", "Verdict"],
    [
        ["gemma3:4b (Google)", "83", "3.3 GB", "✅ Yes", "✅ SELECTED"],
        ["LLaVA-1.5 7B", "28", "4.7 GB", "✅ Yes", "❌ 3× worse OCR"],
        ["LLaVA-1.6 8B", "75", "5.5 GB", "✅ Yes", "❌ Larger, slower"],
        ["BLIP-2", "N/A", "990 MB", "❌ No", "❌ Caption only"],
        ["PaliGemma 2 3B", "81", "~6 GB", "❌ No", "❌ No Ollama"],
    ], top=Inches(1.3))

# Slide: Initial Results (failure)
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_box(slide, "Initial Testing: 75% Failure Rate")
add_table_slide(slide, "Initial Testing: 75% Failure Rate",
    ["Metric", "Result", "Impact"],
    [
        ["Overall failure rate", "75% (6/8)", "Most images get wrong captions"],
        ["Hallucination rate", "50% (4/8)", "Model invents false details"],
        ["Critical hallucinations", "37.5%", "Fake receipts, fake prices"],
        ["Correct text extraction", "25% (2/8)", "Only memes/photos pass"],
    ], top=Inches(1.3))

# Slide: Root Cause
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bullet_slide(slide, "Root Cause Analysis", [
    "768px fixed resize → small text becomes unreadable",
    "4B parameter limit → struggles with multi-element reasoning",
    "No dynamic resolution → entire image at one fixed scale",
    "Hallucination under uncertainty → invents plausible details",
    "No context awareness → confuses screenshots with physical displays",
])

# Slide: Research (5 Papers)
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_box(slide, "Academic Research — 5 Papers Reviewed (2023–2025)")
add_table_slide(slide, "Academic Research — 5 Papers Reviewed (2023–2025)",
    ["Paper", "Key Technique", "Deployable?"],
    [
        ["UReader (Tencent)", "Shape-adaptive cropping", "❌ China-origin"],
        ["TextMonkey (USTC)", "Shifted window attention", "❌ China-origin"],
        ["DocOwl 1.5 (Alibaba)", "Structure-aware parsing", "❌ China-origin"],
        ["InternVL2 (Shanghai AI Lab)", "Tile-based processing", "❌ China-origin"],
        ["Qwen2.5-VL (Alibaba)", "Dynamic resolution + multimodal RoPE", "❌ China-origin"],
    ], top=Inches(1.3))

# Slide: Our Strategy
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_box(slide, "Our Strategy: Take the Ideas, Not the Models")
add_body_box(slide,
    "Constraint: Walmart policy blocks China-origin AI models\n"
    "Solution:   Implement their TECHNIQUES as CODE on our compliant model\n\n"
    "┌──────────────────────────────────────────────────────────┐\n"
    "│ PAPERS SAY:  Use native tiling + dynamic resolution      │\n"
    "│ POLICY SAYS: You can't use those models                  │\n"
    "│ WE DID:     Implement tiling as a multi-pass wrapper     │\n"
    "│             around gemma3:4b (Google — policy compliant)  │\n"
    "└──────────────────────────────────────────────────────────┘\n\n"
    "Technique Mapping:\n"
    "  InternVL2 tile attention  →  Split image into 2-4 crops ourselves\n"
    "  DocOwl structure parsing  →  Pass 1 classifies image type first\n"
    "  UReader focused extraction→  'Read ALL text verbatim' per tile\n"
    "  Qwen2.5-VL resolution    →  Tiling = 2-4× effective resolution",
    top=Inches(1.3), font_size=Pt(13))

# Slide: Multi-Pass Pipeline
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_box(slide, "Multi-Pass Captioning Pipeline (The Fix)")
add_body_box(slide,
    "Pass 1: STRUCTURE  (inspired by DocOwl 1.5)\n"
    "  Full image → 'What TYPE is this?' → screenshot/photo/receipt/meme\n\n"
    "Pass 2: TILE  (inspired by InternVL2)\n"
    "  Split into 2-4 crops → each tile = 2-4× higher effective resolution\n\n"
    "Pass 3: EXTRACT  (inspired by UReader)\n"
    "  Per-tile → 'Read ALL text in this region verbatim'\n\n"
    "Pass 4: MERGE  (text-only LLM call — NO image!)\n"
    "  Combine extracted text → final caption\n"
    "  ⚡ KEY: Model never sees the image in this step\n"
    "         → CANNOT hallucinate visual details\n\n"
    "Same model (gemma3:4b) • Same hardware • Smarter calling strategy",
    top=Inches(1.3), font_size=Pt(14))

# Slide: Results After Fix
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_box(slide, "Results: Hallucination Eliminated (8 + 25 image validation)")
add_table_slide(slide, "Results: Hallucination Eliminated",
    ["Metric", "Before (Single-Pass)", "After (Multi-Pass)", "Change"],
    [
        ["Hallucination rate", "50% (4/8)", "0% (0/8)", "↓ 100%"],
        ["Overall failure rate", "75% (6/8)", "25% (2/8)", "↓ 67%"],
        ["Correct text extraction", "25% (2/8)", "75% (6/8)", "3× better"],
        ["Fabricated claims", "8 total", "0", "Eliminated"],
        ["Avg latency/image", "~5s", "~15s", "3× (acceptable)"],
    ], top=Inches(1.3))

# Slide: 25-Image Scaled Validation
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_box(slide, "Scaled Validation: 25 Images — Confirming at Scale")
add_table_slide(slide, "Scaled Validation: 25 Images",
    ["Verdict", "Count", "Percentage"],
    [
        ["✅ PASS (correct, no hallucination)", "22 / 25", "88%"],
        ["⚠️ PARTIAL (correct but sparse)", "3 / 25", "12%"],
        ["❌ FAIL (missed critical info)", "0 / 25", "0%"],
    ], top=Inches(1.3))
add_body_box(slide,
    "\n\n\n\n\nKey finding: 80% of Walmart Reddit complaint images are screenshots/app screens\n"
    "— exactly the category where single-pass hallucinates most.\n\n"
    "Single-pass hallucination on 25 images: 44% (11/25 images had fabricated details)\n"
    "Multi-pass hallucination on 25 images: 0%",
    top=Inches(3.5), font_size=Pt(13))

# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 3: TRUST SCORE & CONFIDENCE
# ═══════════════════════════════════════════════════════════════════════════════
add_section_divider(prs, 3, "Trust Score & Confidence",
                    "How we validate post credibility and model certainty")

# Slide: Trust Score Formula
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_box(slide, "Trust Score — Weighted Combination Formula")
add_body_box(slide,
    "Formula:\n"
    "  trust_score = 0.4 × metadata + 0.3 × dedup + 0.3 × llm_credibility\n\n"
    "Each component is normalized to [0, 1] and the final score is clamped.\n\n"
    "┌─────────────────────────────────────────────────────────────────┐\n"
    "│ Component 1: METADATA HEURISTICS (weight: 0.4)                  │\n"
    "│   base_floor + 0.20×age + 0.20×karma + 0.30×length + 0.15×eng  │\n"
    "│   • age_score = min(account_age_days / 365, 1.0)                │\n"
    "│   • karma_score = min(total_karma / 5000, 1.0)                  │\n"
    "│   • length_score = min((len(title)+len(body)) / 200, 1.0)       │\n"
    "│   • engagement = min(max(reddit_score, 0) / 20, 1.0)            │\n"
    "├─────────────────────────────────────────────────────────────────┤\n"
    "│ Component 2: DEDUP / ORIGINALITY (weight: 0.3)                  │\n"
    "│   MiniLM-L6-v2 embeddings → cosine similarity                  │\n"
    "│   If max_sim > 0.92 with any previous post → penalize           │\n"
    "├─────────────────────────────────────────────────────────────────┤\n"
    "│ Component 3: LLM CREDIBILITY (weight: 0.3)                      │\n"
    "│   Only invoked when 0.3 < metadata_score < 0.8 (ambiguous zone) │\n"
    "│   Rule-based heuristic (free) or cloud LLM call                │\n"
    "│   Checks: promo language, URL stuffing, caps, retail insider    │\n"
    "└─────────────────────────────────────────────────────────────────┘\n\n"
    "Low-trust posts are FLAGGED for review — never dropped.",
    top=Inches(1.3), font_size=Pt(12))

# Slide: Credibility Signals
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_box(slide, "Credibility Scoring — Negative & Positive Signals")
add_table_slide(slide, "Credibility Scoring — Negative & Positive Signals",
    ["Signal Type", "Indicator", "Score Impact"],
    [
        ["❌ Negative", "Promotional language (≥2 phrases)", "-0.25"],
        ["❌ Negative", "URL stuffing (≥3 links, short text)", "-0.20"],
        ["❌ Negative", "Karma/age mismatch (new acct, high karma)", "-0.20"],
        ["❌ Negative", "Excessive CAPS (>40% letters)", "-0.15"],
        ["❌ Negative", "New account + promotional", "-0.20"],
        ["✅ Positive", "Retail-specific terms (≥2: OGP, ASM, CAP2, etc.)", "+0.25"],
        ["✅ Positive", "Long-form organic text (>600 chars, no links)", "+0.15"],
        ["✅ Positive", "Retail insider terminology (1 term)", "+0.10"],
    ], top=Inches(1.3))

# Slide: Confidence Score
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_box(slide, "Confidence Score — Model Certainty")
add_body_box(slide,
    "Confidence = Softmax probability of the predicted class\n\n"
    "How it's calculated:\n"
    "  1. ModernBERT outputs logits for [negative, neutral, positive]\n"
    "  2. Softmax converts to probabilities: P(neg), P(neu), P(pos)\n"
    "  3. confidence = max(P(neg), P(neu), P(pos))\n\n"
    "Thresholds used in the system:\n"
    "  • Analysis threshold: confidence ≥ 0.7 (from config/models.yaml)\n"
    "  • Notification P1:    trust ≥ 0.70 AND confidence ≥ 0.80\n"
    "  • Notification P2:    trust ≥ 0.50 AND confidence ≥ 0.60\n\n"
    "Combined Priority Formula:\n"
    "  ┌──────────────────────────────────────────────────────┐\n"
    "  │ P1 = (trust_score ≥ 0.70) ∧ (confidence ≥ 0.80)     │\n"
    "  │ P2 = (trust_score ≥ 0.50) ∧ (confidence ≥ 0.60)     │\n"
    "  │ Below P2 thresholds → no notification triggered       │\n"
    "  └──────────────────────────────────────────────────────┘\n\n"
    "Dashboard displays both scores per post for analyst transparency.",
    top=Inches(1.3), font_size=Pt(13))

# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 4: MODERNBERT
# ═══════════════════════════════════════════════════════════════════════════════
add_section_divider(prs, 4, "ModernBERT Fine-Tuning",
                    "Domain-specialized long-context sentiment classification")

# Slide: Why ModernBERT
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_box(slide, "Why ModernBERT over RoBERTa?")
add_body_box(slide,
    "Thesis Claim:\n"
    "  Long-context, domain-fine-tuned encoders beat short-context\n"
    "  Twitter-trained baselines on Reddit-flavored retail complaints.\n\n"
    "RoBERTa Limitations:\n"
    "  • 512 token context window — truncates long Reddit posts\n"
    "  • Trained on Twitter data — different register from Reddit\n"
    "  • No domain knowledge of retail/Walmart terminology\n\n"
    "ModernBERT Advantages:\n"
    "  • 8192 token context window (16× longer)\n"
    "  • Supports full Reddit complaint posts without truncation\n"
    "  • Fine-tunable with curriculum learning\n"
    "  • Modern architecture optimizations",
    top=Inches(1.3), font_size=Pt(14))

# Slide: 3-Stage Curriculum
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_box(slide, "3-Stage Curriculum Training Pipeline")
add_table_slide(slide, "3-Stage Curriculum Training Pipeline",
    ["Stage", "Dataset", "Epochs", "Purpose"],
    [
        ["1. Generic Sentiment", "TweetEval-sentiment (45K tweets)", "2", "Polarity grounding"],
        ["2. Reddit Register", "GoEmotions-3class (54K Reddit comments)", "2", "Reddit language patterns"],
        ["3. Domain Specialization", "Walmart-200 (5-fold CV)", "up to 15 (patience 3)", "Retail-specific fine-tuning"],
    ], top=Inches(1.3))
add_body_box(slide,
    "\n\n\n\n\nTraining Configuration:\n"
    "  • max_length: 1024 tokens (key lever for long-context advantage)\n"
    "  • Effective batch size: 32 (per-device BS=8 × grad accum=4)\n"
    "  • Class weights: neg=0.52, neu=1.03, pos=8.33 (inverse frequency)\n"
    "  • Minority oversampling to ~100/class per training fold\n"
    "  • Early stopping on eval_macro_f1 with patience=3\n"
    "  • Hardware: Apple M-series (MPS backend)",
    top=Inches(3.5), font_size=Pt(13))

# Slide: Challenges Faced
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_box(slide, "Key Challenges & Mitigations")
add_table_slide(slide, "Key Challenges & Mitigations",
    ["Challenge", "Impact", "Mitigation"],
    [
        ["Synthetic benchmark (77-char bodies)", "Cannot test long-context", "Built real 200-post benchmark (min 300 chars)"],
        ["Corp network blocked HuggingFace", "Cannot download models", "Restart + hotspot + offline triad env vars"],
        ["Stage 3 macro F1 = 0.40 (catastrophic)", "Positive class collapsed", "Class weights + oversampling + curriculum"],
        ["Eval showed F1=1.0 (memorization)", "Leakage — useless result", "Switched to out-of-fold CV predictions"],
        ["Long-context not showing up", "max_length=512 (defeats purpose)", "Retrained at max_length=1024"],
        ["AI-assist acceptance = 100%", "Defensibility concern", "Disclosed + blind recheck planned"],
    ], top=Inches(1.3))

# Slide: Final Results
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_box(slide, "Final Results — ModernBERT vs RoBERTa")
add_table_slide(slide, "Final Results — ModernBERT vs RoBERTa (Out-of-Fold CV)",
    ["Metric", "RoBERTa", "ModernBERT v2", "Improvement"],
    [
        ["Macro F1 (overall)", "0.6272", "0.7642", "+0.137 (+22%)"],
        ["F1 negative", "0.7967", "0.8779", "+0.081"],
        ["F1 neutral", "0.6087", "0.7480", "+0.139"],
        ["F1 positive", "0.4762", "0.6667", "+0.190"],
        ["Long-post F1 (≥512 tokens)", "0.2778", "1.0000", "+0.722 (✨)"],
        ["Short-post F1 (n=193)", "0.6360", "0.7619", "+0.126"],
        ["Latency (ms/post, MPS)", "6.5 ms", "11.9 ms", "+5.4 ms"],
    ], top=Inches(1.3))

# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 5: DASHBOARD
# ═══════════════════════════════════════════════════════════════════════════════
add_section_divider(prs, 5, "Dashboard & Data Population",
                    "From pipeline output to actionable insights on screen")

# Slide: Dashboard Overview
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_box(slide, "Dashboard Pages — Navigation Structure")
add_table_slide(slide, "Dashboard Pages — Navigation Structure",
    ["Page", "Priority", "Purpose", "Key Metric"],
    [
        ["Brand Health", "P0", "At-a-glance KPIs & trends", "Overall sentiment score"],
        ["Post Explorer", "P1", "Search/filter all analyzed posts", "Volume + sentiment dist"],
        ["Review & Validate", "P0", "Correct labels + draft replies", "Accuracy improvement"],
        ["Post Lifecycle", "P0", "Kanban workflow (triaged→resolved)", "Resolution rate"],
        ["Insights", "P1", "AI-generated competitor analysis", "Issue rankings"],
        ["Pipeline Control", "P1", "Monitor & trigger runs", "Jobs, cursors, health"],
        ["Notifications", "P1", "Group-based alert routing config", "Delivery log"],
    ], top=Inches(1.3))

# Slide: Brand Health KPIs
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_box(slide, "Brand Health — KPI Tiles with Breakdowns")
add_body_box(slide,
    "5 KPI Tiles (each with embedded breakdown + dual navigation):\n\n"
    "┌──────────┐  ┌──────────┐  ┌──────────┐  ┌──────────┐  ┌──────────┐\n"
    "│  Total   │  │ Negative │  │ Priority │  │ Positive │  │Lifecycle │\n"
    "│  Posts   │  │  Posts   │  │  (P1+P2) │  │  Posts   │  │  Status  │\n"
    "│  ────    │  │  ────    │  │  ────    │  │  ────    │  │  ────    │\n"
    "│ breakdown│  │ by conf  │  │ P1 / P2  │  │ by conf  │  │ by state │\n"
    "└──────────┘  └──────────┘  └──────────┘  └──────────┘  └──────────┘\n\n"
    "Each tile supports:\n"
    "  • Single-click → navigate to Post Explorer (filtered)\n"
    "  • Dropdown → dual nav to Post Explorer OR Review & Reply\n\n"
    "Data populated from:\n"
    "  • /api/brand-health?range=today (aggregated sentiment counts)\n"
    "  • /api/brand-health/priority-negatives (P1+P2 tier counts)\n"
    "  • /api/lifecycle (state distribution)",
    top=Inches(1.3), font_size=Pt(13))

# Slide: Data Flow to Dashboard
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_box(slide, "How Data Populates the Dashboard")
add_body_box(slide,
    "Pipeline Output → SQLite Tables → FastAPI Endpoints → React Components\n\n"
    "Mapping:\n"
    "  ┌─────────────────┐    ┌────────────────────────┐    ┌───────────────┐\n"
    "  │ analyses table   │ →  │ /api/brand-health       │ →  │ KPI Tiles     │\n"
    "  │ (sentiment, conf)│    │ /api/posts              │    │ Trend Charts  │\n"
    "  │ (aspects, trust) │    │ /api/segments           │    │ Aspect Heatmap│\n"
    "  ├─────────────────┤    ├────────────────────────┤    ├───────────────┤\n"
    "  │ aggregates table │ →  │ /api/brand-health       │ →  │ Trend Line    │\n"
    "  │ (hourly/daily)   │    │   ?range=week/month     │    │ Volume Ticker │\n"
    "  ├─────────────────┤    ├────────────────────────┤    ├───────────────┤\n"
    "  │ alerts table     │ →  │ /api/alerts             │ →  │ Alert Feed    │\n"
    "  │ (spike, severity)│    │ WebSocket push          │    │ Real-time     │\n"
    "  ├─────────────────┤    ├────────────────────────┤    ├───────────────┤\n"
    "  │ feedback table   │ →  │ /api/review/*           │ →  │ Review Queue  │\n"
    "  │ (corrections)    │    │ /api/review/{id}/draft  │    │ Draft Replies │\n"
    "  └─────────────────┘    └────────────────────────┘    └───────────────┘\n\n"
    "WebSocket for real-time alert push — no polling needed.",
    top=Inches(1.3), font_size=Pt(12))

# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 6: NOTIFICATIONS
# ═══════════════════════════════════════════════════════════════════════════════
add_section_divider(prs, 6, "Notification System",
                    "Group-based routing for P1/P2 priority posts")

# Slide: Notification Architecture
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_box(slide, "Notification System — Group-Based Routing")
add_body_box(slide,
    "Architecture:\n\n"
    "  Pipeline analyzes post → classify_priority(trust, confidence)\n"
    "      │\n"
    "      ├── Not P1/P2? → Skip (no notification)\n"
    "      │\n"
    "      └── P1 or P2? → Find matching notification groups\n"
    "              │\n"
    "              ├── Group matches subreddit + priority filter?\n"
    "              │       │\n"
    "              │       ├── Email DL configured → Send email\n"
    "              │       └── Slack channel configured → Send Slack\n"
    "              │\n"
    "              └── Log to notification_log table (audit trail)\n\n"
    "Priority Classification:\n"
    "  P1: trust ≥ 0.70 AND confidence ≥ 0.80 (high-signal, immediate action)\n"
    "  P2: trust ≥ 0.50 AND confidence ≥ 0.60 (review-worthy, lower urgency)\n\n"
    "Groups are configured per subreddit-set → different teams own different subreddits",
    top=Inches(1.3), font_size=Pt(13))

# Slide: Notification Config Page
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_box(slide, "Notification Configuration Page (Frontend)")
add_body_box(slide,
    "Admin UI at /notifications allows:\n\n"
    "  • Create notification groups (name, subreddits, email DL, Slack)\n"
    "  • Quick-add subreddits by category (Walmart core, Spark, Pharmacy...)\n"
    "  • Priority filter: choose P1, P2, or both per group\n"
    "  • Enable/disable groups with toggle\n"
    "  • Test (dry-run) — simulates notification without sending\n"
    "  • View delivery log — audit trail of all sent notifications\n"
    "  • Delete groups\n\n"
    "API Endpoints (8 total):\n"
    "  GET  /api/notifications/config          — overall config + groups\n"
    "  GET  /api/notifications/groups          — list all groups\n"
    "  POST /api/notifications/groups          — create new group\n"
    "  PUT  /api/notifications/groups/{id}     — update group\n"
    "  DELETE /api/notifications/groups/{id}   — delete group\n"
    "  POST /api/notifications/test/{id}       — dry-run test\n"
    "  GET  /api/notifications/log             — delivery audit\n"
    "  GET  /api/notifications/subreddits      — available subreddit list",
    top=Inches(1.3), font_size=Pt(12))

# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 7: REVIEW & VALIDATE
# ═══════════════════════════════════════════════════════════════════════════════
add_section_divider(prs, 7, "Review & Validate",
                    "Human-in-the-loop correction and reply generation")

# Slide: Review Queue
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_box(slide, "Review & Validate — Human-in-the-Loop")
add_body_box(slide,
    "Purpose: Analysts correct AI labels and draft replies to negative posts\n\n"
    "Workflow:\n"
    "  1. Queue shows posts sorted by priority (P1 first)\n"
    "  2. Analyst reviews sentiment + aspects (correct if wrong)\n"
    "  3. Click 'Generate Drafts' → two reply options generated\n\n"
    "Reply Generation (Dual-Draft Picker):\n"
    "  ┌─────────────────────┐    ┌─────────────────────┐\n"
    "  │  Draft A             │    │  Draft B             │\n"
    "  │  Smart Composer      │    │  FLAN-T5-base        │\n"
    "  │  (keyword extraction │    │  (multi-temp sampling │\n"
    "  │   + phrase pools)    │    │   + scorer)           │\n"
    "  └─────────────────────┘    └─────────────────────┘\n\n"
    "  Analyst picks one → edits → posts to Reddit\n\n"
    "Learning Loop:\n"
    "  Posted replies saved to feedback table → become few-shot examples\n"
    "  for future draft generation (tone matching improves over time)",
    top=Inches(1.3), font_size=Pt(13))

# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 8: POST EXPLORER
# ═══════════════════════════════════════════════════════════════════════════════
add_section_divider(prs, 8, "Post Explorer & Filtering",
                    "Search, filter, and deep-dive into analyzed posts")

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_box(slide, "Post Explorer — Features")
add_body_box(slide,
    "Purpose: Browse all analyzed posts with powerful filtering\n\n"
    "Filters Available:\n"
    "  • Sentiment: negative / neutral / positive\n"
    "  • Confidence: threshold slider\n"
    "  • Trust score: threshold slider\n"
    "  • Subreddit: multi-select from 25 tracked\n"
    "  • Aspect: delivery, product_quality, returns, support, pricing, app\n"
    "  • Date range: today, week, month, custom\n"
    "  • Text search: full-text search across titles & bodies\n\n"
    "Per-Post Card Shows:\n"
    "  • Title + body excerpt\n"
    "  • Sentiment badge (color-coded) + confidence %\n"
    "  • Trust score indicator\n"
    "  • Aspect tags\n"
    "  • Subreddit + post time\n"
    "  • Reddit link → open original thread\n"
    "  • Actions: Review, Add to Lifecycle, View Details",
    top=Inches(1.3), font_size=Pt(13))

# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 9: POST LIFECYCLE
# ═══════════════════════════════════════════════════════════════════════════════
add_section_divider(prs, 9, "Post Lifecycle",
                    "Kanban workflow from triage to resolution")

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_box(slide, "Post Lifecycle — Kanban Board")
add_body_box(slide,
    "States: Triaged → Acknowledged → In Progress → Resolved\n\n"
    "┌─────────────┐  ┌─────────────┐  ┌─────────────┐  ┌─────────────┐\n"
    "│   TRIAGED   │  │ ACKNOWLEDGED│  │ IN PROGRESS │  │  RESOLVED   │\n"
    "│             │  │             │  │             │  │             │\n"
    "│  New P1/P2  │  │  Assigned   │  │  Reply being│  │  Reply      │\n"
    "│  posts land │  │  to analyst │  │  drafted    │  │  posted     │\n"
    "│  here       │  │             │  │             │  │             │\n"
    "└──────┬──────┘  └──────┬──────┘  └──────┬──────┘  └─────────────┘\n"
    "       │                │                │\n"
    "       └── Acknowledge ─┘── Start Work ──┘── Resolve\n\n"
    "Resolve Modal (2-step flow):\n"
    "  Step 1: Save action note + optional LLM-drafted reply\n"
    "          → 'Save reply & open Reddit' OR 'Resolve (no reply needed)'\n"
    "  Step 2: Copy reply to clipboard → Open Reddit thread → Mark Resolved\n\n"
    "All state transitions are logged with timestamps for SLA tracking.",
    top=Inches(1.3), font_size=Pt(13))

# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 10: INSIGHTS
# ═══════════════════════════════════════════════════════════════════════════════
add_section_divider(prs, 10, "Insights & Competitor Analysis",
                    "AI-generated strategic intelligence from raw data")

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_box(slide, "Insights Page — What It Provides")
add_body_box(slide,
    "AI-Generated Insights:\n\n"
    "  1. Issue Rankings\n"
    "     • Top negative issues ranked by volume × severity × recency\n"
    "     • Grouped by aspect (delivery, pricing, app, etc.)\n"
    "     • Trend arrows showing improvement or deterioration\n\n"
    "  2. Competitor Pulse\n"
    "     • Compare sentiment across Walmart vs competitors\n"
    "     • Tracked competitors: Costco, Target, Amazon\n"
    "     • Cross-mentioned posts (\"Walmart vs Costco\" comparisons)\n\n"
    "  3. LLM Summarization\n"
    "     • Natural language summaries of week's top themes\n"
    "     • Suggested action items for product teams\n"
    "     • Emerging topic detection (new phrase clusters)\n\n"
    "  4. Aspect Drilldown\n"
    "     • 6-category taxonomy: delivery, product_quality, returns,\n"
    "       customer_support, pricing, app_website\n"
    "     • Per-aspect sentiment trends, volume, representative posts",
    top=Inches(1.3), font_size=Pt(13))

# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 11: RESULTS & FUTURE WORK
# ═══════════════════════════════════════════════════════════════════════════════
add_section_divider(prs, 11, "Results & Future Work",
                    "What we achieved and what's next")

# Slide: Key Achievements
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_box(slide, "Key Achievements — Post Midsem")
add_table_slide(slide, "Key Achievements — Post Midsem",
    ["Area", "Achievement", "Evidence"],
    [
        ["ModernBERT", "Macro F1: 0.6272 → 0.7642 (+22%)", "5-fold OOF CV"],
        ["Long-context", "Long-post F1: 0.28 → 1.00 (+722%)", "7 posts ≥512 tokens"],
        ["Vision", "Hallucination: 50% → 0%", "8+25 image validation"],
        ["Vision", "Text extraction: 25% → 75%", "Multi-pass pipeline"],
        ["Pipeline", "25 subreddits, hourly automated", "Arctic Shift API"],
        ["Dashboard", "7 pages, real-time WebSocket alerts", "React + FastAPI"],
        ["Notifications", "Group-based P1/P2 routing", "Email + Slack channels"],
        ["Lifecycle", "Full Kanban with 2-step resolve", "Triage → Resolved"],
    ], top=Inches(1.3))

# Slide: Technical Stack Summary
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_box(slide, "Technical Stack Summary")
add_table_slide(slide, "Technical Stack Summary",
    ["Layer", "Technology", "Key Decisions"],
    [
        ["Backend", "Python 3.13 + FastAPI + SQLite", "Free, local-first, modular"],
        ["Frontend", "React 18 + TypeScript + Vite + Tailwind", "Modern SPA, responsive"],
        ["Sentiment", "ModernBERT (fine-tuned, 1024 tokens)", "Domain-specialized, offline"],
        ["Aspects", "BART-MNLI (zero-shot classification)", "No training needed"],
        ["Vision", "Gemma 3 4B via Ollama (multi-pass)", "Policy compliant, no hallucination"],
        ["Reply Gen", "FLAN-T5 + Smart Composer (dual draft)", "Learning loop via feedback"],
        ["Trust", "Metadata + Dedup + LLM (weighted formula)", "Flag, don't drop"],
        ["Scheduling", "asyncio lifespan (60-min) + manual", "Cursor-based incremental"],
    ], top=Inches(1.3))

# Slide: Future Work
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title_box(slide, "Future Work")
add_bullet_slide(slide, "Future Work", [
    "Azure Cosmos DB migration (production-grade storage)",
    "Twitter/X integration as second data source",
    "3-seed ensemble for ModernBERT (+0.01–0.03 F1, tighter variance)",
    "Gemma 3 12B upgrade for remaining edge-case images",
    "Spanish language support (bilingual retail communities)",
    "Azure AD authentication for multi-user access",
    "Automated retraining pipeline (feedback loop → model updates)",
    "Slack bot integration for inline notification responses",
])

# ─── SLIDE: Thank You ────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BLUE)
add_title_box(slide, "Thank You",
              top=Inches(2.5), left=Inches(0.8), font_size=Pt(40),
              color=WHITE)
add_body_box(slide,
             "Retail Sentiment Intelligence\n"
             "Post Mid-Semester Presentation\n\n"
             "Questions?",
             top=Inches(3.8), left=Inches(0.8), font_size=Pt(20),
             color=RGBColor(0xBB, 0xBB, 0xBB))

# ═══════════════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════════════
os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
prs.save(OUTPUT_PATH)
print(f"✅ Presentation saved: {OUTPUT_PATH}")
print(f"   Slides: {len(prs.slides)}")
