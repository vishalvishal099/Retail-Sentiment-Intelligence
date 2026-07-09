"""
Generate the Post-Midsem slide deck (v2) aligned with the current
docs/Sem_4/RSI_Post_Midsem_Presentation.md content:

  - BITS Pilani WILP affiliation (not LPU)
  - Every-6h scheduler cadence (not 60 min)
  - Part A (mid-sem recap, 6 topics) / Part B (post-midsem, 5 topics) split
  - New ModernBERT "Training Evidence" slide
  - Vision multi-pass story + ModernBERT curriculum both retained

Output:  docs/Sem_4/RSI_Post_Midsem_Presentation_v2.pptx
"""

from __future__ import annotations
import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Paths & style ───────────────────────────────────────────────────────────
OUT = Path(__file__).parent / "Sem_4" / "RSI_Post_Midsem_Presentation_v2.pptx"

# Brand palette
WALMART_BLUE = RGBColor(0x00, 0x71, 0xDC)
DARK_BLUE    = RGBColor(0x04, 0x1E, 0x42)
LIGHT_BLUE   = RGBColor(0xE8, 0xF4, 0xFD)
YELLOW       = RGBColor(0xFF, 0xC2, 0x20)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY    = RGBColor(0x33, 0x33, 0x33)
MED_GRAY     = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY   = RGBColor(0xF5, 0xF5, 0xF5)
GREEN        = RGBColor(0x10, 0xB9, 0x81)
RED          = RGBColor(0xEF, 0x44, 0x44)
PURPLE       = RGBColor(0x7C, 0x3A, 0xED)
AMBER        = RGBColor(0xD9, 0x77, 0x06)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ─── Helpers ─────────────────────────────────────────────────────────────────
def _set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _tx(slide, left, top, width, height, text, *,
        size=14, bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        for r in p.runs:
            r.font.size = size
            r.font.bold = bold
            r.font.color.rgb = color
    return box


def _rect(slide, left, top, width, height, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def _header_bar(slide, title, subtitle=None):
    """Blue header bar with white title."""
    _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.85), DARK_BLUE)
    _tx(slide, Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.5),
        title, size=Pt(26), bold=True, color=WHITE)
    if subtitle:
        _tx(slide, Inches(0.4), Inches(0.5), Inches(12.5), Inches(0.35),
            subtitle, size=Pt(12), color=LIGHT_BLUE)
    # thin accent line
    _rect(slide, Inches(0), Inches(0.85), SLIDE_W, Inches(0.05), YELLOW)


def _footer(slide, page_no, total):
    _tx(slide, Inches(0.4), Inches(7.15), Inches(6), Inches(0.3),
        "Retail Sentiment Intelligence  ·  BITS ZG628T Dissertation  ·  Vishal Singh",
        size=Pt(9), color=MED_GRAY)
    _tx(slide, Inches(12.0), Inches(7.15), Inches(1.2), Inches(0.3),
        f"{page_no} / {total}", size=Pt(9), color=MED_GRAY, align=PP_ALIGN.RIGHT)


def _bullets(slide, left, top, width, height, items, *,
             size=Pt(14), color=DARK_GRAY, bold=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.space_after = Pt(4)
        for r in p.runs:
            r.font.size = size
            r.font.color.rgb = color
            r.font.bold = bold
    return box


def _table(slide, left, top, width, height, data, *,
           header_fill=DARK_BLUE, header_color=WHITE,
           body_font=Pt(11), first_col_bold=False,
           highlight_rows=None):
    """data = list of rows, first row is header."""
    highlight_rows = highlight_rows or {}
    rows = len(data)
    cols = len(data[0])
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = tbl_shape.table
    for r_idx, row in enumerate(data):
        for c_idx, val in enumerate(row):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = str(val)
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
            elif r_idx in highlight_rows:
                cell.fill.solid()
                cell.fill.fore_color.rgb = highlight_rows[r_idx]
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r_idx % 2 else LIGHT_GRAY
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = body_font
                    if r_idx == 0:
                        run.font.color.rgb = header_color
                        run.font.bold = True
                    else:
                        run.font.color.rgb = DARK_GRAY
                        if first_col_bold and c_idx == 0:
                            run.font.bold = True
    return tbl


def _kpi_tile(slide, left, top, width, height, label, value, *,
              tint=LIGHT_BLUE, value_color=WALMART_BLUE):
    _rect(slide, left, top, width, height, tint, line=WALMART_BLUE)
    _tx(slide, left, top + Inches(0.15), width, Inches(0.35),
        label, size=Pt(11), bold=True, color=MED_GRAY, align=PP_ALIGN.CENTER)
    _tx(slide, left, top + Inches(0.5), width, Inches(0.9),
        value, size=Pt(28), bold=True, color=value_color, align=PP_ALIGN.CENTER)


def _pipeline_stage(slide, left, top, width, height, num, title, subtitle):
    _rect(slide, left, top, width, height, LIGHT_BLUE, line=WALMART_BLUE)
    _tx(slide, left, top + Inches(0.1), width, Inches(0.3),
        f"STAGE {num}", size=Pt(9), bold=True, color=WALMART_BLUE,
        align=PP_ALIGN.CENTER)
    _tx(slide, left, top + Inches(0.4), width, Inches(0.35),
        title, size=Pt(13), bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    _tx(slide, left, top + Inches(0.75), width, Inches(0.7),
        subtitle, size=Pt(9), color=MED_GRAY, align=PP_ALIGN.CENTER)


def _arrow(slide, x, y, w=Inches(0.35), h=Inches(0.35)):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    a.fill.solid()
    a.fill.fore_color.rgb = YELLOW
    a.line.fill.background()


# ─── Slide builders ──────────────────────────────────────────────────────────
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total = 32  # keep in sync with # of add_* calls below
    page = [0]  # mutable counter

    def new():
        page[0] += 1
        s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        return s

    # ── 1. Title ────────────────────────────────────────────────────────────
    s = new()
    _set_bg(s, DARK_BLUE)
    # yellow accent bar
    _rect(s, Inches(0), Inches(2.6), Inches(0.4), Inches(2.5), YELLOW)
    _tx(s, Inches(0.9), Inches(1.4), Inches(11.5), Inches(0.9),
        "Retail Sentiment Intelligence", size=Pt(46), bold=True, color=WHITE)
    _tx(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(0.6),
        "Real-Time Brand Health Monitoring via Reddit NLP Pipeline",
        size=Pt(22), color=YELLOW)
    _tx(s, Inches(0.9), Inches(3.1), Inches(11.5), Inches(0.4),
        "BITS ZG628T  ·  Dissertation  ·  Post Mid-Semester Progress Presentation",
        size=Pt(14), bold=True, color=WHITE)
    _tx(s, Inches(0.9), Inches(4.1), Inches(11.5), Inches(0.4),
        "Vishal Singh   ·   ID No. 2020AA05641",
        size=Pt(18), bold=True, color=WHITE)
    _tx(s, Inches(0.9), Inches(4.55), Inches(11.5), Inches(0.4),
        "M.Tech in Artificial Intelligence & Machine Learning",
        size=Pt(14), color=LIGHT_BLUE)
    _tx(s, Inches(0.9), Inches(4.9), Inches(11.5), Inches(0.4),
        "Birla Institute of Technology & Science, Pilani (WILP)",
        size=Pt(14), color=LIGHT_BLUE)
    _tx(s, Inches(0.9), Inches(5.6), Inches(11.5), Inches(0.4),
        "Dissertation work carried out at  Walmart Global Tech, Bengaluru",
        size=Pt(12), color=WHITE)
    _tx(s, Inches(0.9), Inches(5.95), Inches(11.5), Inches(0.4),
        "Supervisor:  Mr. Varunendra Pratap Singh — Principal Software Engineer, Walmart Global Tech",
        size=Pt(12), color=WHITE)
    _tx(s, Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.4),
        "July 2026", size=Pt(11), color=MED_GRAY)

    # ── 2. Agenda ───────────────────────────────────────────────────────────
    s = new()
    _header_bar(s, "Agenda", "Part A recap · Part B is the focus of today")
    # Part A
    _rect(s, Inches(0.4), Inches(1.2), Inches(6.2), Inches(5.5), LIGHT_GRAY,
          line=MED_GRAY)
    _tx(s, Inches(0.6), Inches(1.35), Inches(6.0), Inches(0.4),
        "PART A  ·  MID-SEM RECAP  (already presented)",
        size=Pt(13), bold=True, color=AMBER)
    _bullets(s, Inches(0.7), Inches(1.85), Inches(5.9), Inches(4.7), [
        "System Architecture & End-to-End Pipeline",
        "Vision / Image Processing — Challenge & Mitigation",
        "Trust Score & Confidence Calculations",
        "ModernBERT — Domain Fine-Tuning Journey",
        "Dashboard — Data Population & Sections",
        "Notification System — Group-Based Routing",
    ], size=Pt(15), color=DARK_GRAY)
    # Part B
    _rect(s, Inches(6.85), Inches(1.2), Inches(6.1), Inches(5.5), LIGHT_BLUE,
          line=WALMART_BLUE)
    _tx(s, Inches(7.05), Inches(1.35), Inches(5.9), Inches(0.4),
        "PART B  ·  POST-MIDSEM WORK  (new — today's focus)",
        size=Pt(13), bold=True, color=WALMART_BLUE)
    _bullets(s, Inches(7.15), Inches(1.85), Inches(5.85), Inches(4.7), [
        "Review & Validate — Human-in-the-Loop",
        "Post Explorer & Filtering",
        "Post Lifecycle — Kanban Workflow",
        "Insights & Competitor Analysis",
        "Results & Future Work",
    ], size=Pt(15), color=DARK_GRAY, bold=True)
    _footer(s, page[0], total)

    # ── 3. Part A divider ────────────────────────────────────────────────────
    s = new()
    _set_bg(s, LIGHT_GRAY)
    _rect(s, Inches(0), Inches(3.0), SLIDE_W, Inches(1.5), AMBER)
    _tx(s, Inches(0.5), Inches(3.15), Inches(12.3), Inches(0.7),
        "PART A", size=Pt(30), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _tx(s, Inches(0.5), Inches(3.65), Inches(12.3), Inches(0.7),
        "Mid-Semester Recap", size=Pt(40), bold=True, color=WHITE,
        align=PP_ALIGN.CENTER)
    _tx(s, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.5),
        "6 topics · ~8 minutes · quick walkthrough of what was built and demoed at mid-sem",
        size=Pt(14), color=DARK_GRAY, align=PP_ALIGN.CENTER)
    _footer(s, page[0], total)

    # ── 4. Recap 1 · High-Level Architecture (layered) ───────────────────────
    s = new()
    _header_bar(s, "Recap 1 · System Architecture (layered view)",
                "Sources → Pipeline → AI Runtime → Storage → Serving → Clients")
    layers = [
        ("L1  Data Sources",        "Arctic Shift API  ·  PRAW (optional)",             AMBER),
        ("L2  Pipeline",            "Ingest · Preprocess · Trust · Analyze · Aggregate · Alert",  WALMART_BLUE),
        ("L3  AI / ML Runtime",     "ModernBERT · BART-MNLI · Gemma 3 4B (Ollama) · FLAN-T5", PURPLE),
        ("L4  Storage",             "SQLite (dev)  /  Azure Cosmos DB (prod)  —  6 containers", RED),
        ("L5  Serving",             "FastAPI  ·  REST + WebSocket  ·  port 8001",       GREEN),
        ("L6  Clients",             "React Dashboard  ·  Email DL  ·  Slack",           DARK_BLUE),
    ]
    y = Inches(1.15)
    for name, desc, col in layers:
        _rect(s, Inches(0.5), y, Inches(12.3), Inches(0.85), LIGHT_GRAY,
              line=col)
        _rect(s, Inches(0.5), y, Inches(0.25), Inches(0.85), col)  # left stripe
        _tx(s, Inches(0.9), y + Inches(0.05), Inches(3.5), Inches(0.4),
            name, size=Pt(14), bold=True, color=col)
        _tx(s, Inches(0.9), y + Inches(0.42), Inches(11.5), Inches(0.4),
            desc, size=Pt(12), color=DARK_GRAY)
        y = y + Inches(0.95)
    _footer(s, page[0], total)

    # ── 5. Recap 1 · Component Stack table ───────────────────────────────────
    s = new()
    _header_bar(s, "Recap 1 · Component Stack Summary")
    _table(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(4.5), [
        ["Layer", "Stack", "Notes"],
        ["Data Sources",   "Arctic Shift API  +  PRAW (optional)",           "Free, no auth needed"],
        ["Pipeline",       "Python 3.13  ·  asyncio scheduler",              "6-hour cadence + on-demand"],
        ["AI Runtime",     "HuggingFace (offline) + Ollama + Azure OpenAI",  "Local-first, modular"],
        ["Storage",        "SQLite (dev)  /  Azure Cosmos DB (prod)",        "Pluggable backend"],
        ["API",            "FastAPI  +  WebSocket",                          "Port 8001"],
        ["Frontend",       "React 18  ·  TypeScript  ·  Vite  ·  Tailwind",  "Port 5173"],
        ["Notifications",  "SMTP  +  Slack Webhook",                         "Group-based routing"],
    ], first_col_bold=True, highlight_rows={2: LIGHT_BLUE})
    _tx(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.4),
        "Highlighted row — corrected from mid-sem deck (was 60 min, is actually every 6 h in production).",
        size=Pt(11), color=AMBER, bold=True)
    _footer(s, page[0], total)

    # ── 6. Recap 1 · 6-stage pipeline visual ────────────────────────────────
    s = new()
    _header_bar(s, "Recap 1 · 6-Stage Pipeline",
                "Scheduler tick every 6 h · or manual Run Now")
    stages = [
        ("1", "INGEST",    "Arctic Shift\n25 subreddits"),
        ("2", "PREPROCESS","Clean · langdetect\nMiniLM dedup > 0.92"),
        ("3", "TRUST",     "meta + dedup + LLM\nweighted 0.4/0.3/0.3"),
        ("4", "ANALYZE",   "ModernBERT · BART-MNLI\nGemma 3 4B (if image)"),
        ("5", "AGGREGATE", "Hourly / daily rollups\nper-subreddit · per-aspect"),
        ("6", "ALERT",     "Spike + severity\nP1/P2 → notify groups"),
    ]
    x = Inches(0.35)
    top = Inches(2.3)
    box_w = Inches(1.85)
    box_h = Inches(1.7)
    gap = Inches(0.15)
    for i, (n, t, sub) in enumerate(stages):
        _pipeline_stage(s, x, top, box_w, box_h, n, t, sub)
        if i < 5:
            _arrow(s, x + box_w + Inches(0.005),
                   top + Inches(0.65), Inches(0.15), Inches(0.35))
        x = x + box_w + gap
    _tx(s, Inches(0.5), Inches(4.4), Inches(12.3), Inches(0.4),
        "Total latency  ~3–5 min for 25 subreddits  ·  cursor-based incremental",
        size=Pt(13), bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    _tx(s, Inches(0.5), Inches(4.9), Inches(12.3), Inches(0.4),
        "→ live alerts pushed to React dashboard via WebSocket",
        size=Pt(12), color=MED_GRAY, align=PP_ALIGN.CENTER)
    _footer(s, page[0], total)

    # ── 7. Recap 1 · Data flow ──────────────────────────────────────────────
    s = new()
    _header_bar(s, "Recap 1 · Data Flow — Raw Post to Dashboard KPI")
    stages = [
        ("RAW POST",   "title · body\nmedia · author",       AMBER),
        ("CLEAN",      "strip HTML\nlang · dedup",           WALMART_BLUE),
        ("SCORE",      "trust · sentiment\naspects · vision",PURPLE),
        ("PERSIST",    "raw_posts · analyses\naggregates · alerts", RED),
        ("SURFACE",    "KPIs · aspect drill\nalert feed · explorer", GREEN),
    ]
    x = Inches(0.4)
    top = Inches(2.4)
    box_w = Inches(2.35)
    box_h = Inches(1.9)
    gap = Inches(0.15)
    for i, (t, sub, col) in enumerate(stages):
        _rect(s, x, top, box_w, box_h, LIGHT_GRAY, line=col)
        _rect(s, x, top, box_w, Inches(0.4), col)
        _tx(s, x, top + Inches(0.05), box_w, Inches(0.35),
            t, size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _tx(s, x, top + Inches(0.6), box_w, Inches(1.2),
            sub, size=Pt(12), color=DARK_GRAY, align=PP_ALIGN.CENTER)
        if i < 4:
            _arrow(s, x + box_w + Inches(0.005),
                   top + Inches(0.75), Inches(0.14), Inches(0.35))
        x = x + box_w + gap
    _tx(s, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.4),
        "Every post is scored, stored, and surfaced — the trust score decides whether an alert fires.",
        size=Pt(12), color=DARK_GRAY, align=PP_ALIGN.CENTER)
    _footer(s, page[0], total)

    # ── 8. Recap 1 · Storage schema ─────────────────────────────────────────
    s = new()
    _header_bar(s, "Recap 1 · Storage Schema")
    _table(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.2), [
        ["Table", "Partition Key", "Key Fields", "Purpose"],
        ["raw_posts",             "/subreddit",     "id, title, body, author_hash, created_utc",  "Ingested (privacy-safe)"],
        ["analyses",              "/subreddit",     "post_id, sentiment, confidence, aspects, trust_score", "AI analysis results"],
        ["aggregates",            "/time_window",   "subreddit, window, metrics_json",            "Pre-computed KPIs"],
        ["feedback",              "/analyst_id",    "post_id, correction, reply_text",            "Human corrections + replies"],
        ["alerts",                "/severity",      "type, aspect, threshold_breached",           "Triggered anomalies"],
        ["notification_groups",   "id",             "subreddits[], email_dl[], priority_filter",  "Routing config"],
        ["notification_log",      "group_id",       "post_id, channel, status, sent_at",          "Delivery audit trail"],
    ], first_col_bold=True)
    _tx(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4),
        "Privacy — Reddit usernames are SHA-hashed before storage · 1-year retention default",
        size=Pt(11), bold=True, color=AMBER)
    _footer(s, page[0], total)

    # ── 9. Recap 2 · Vision problem ─────────────────────────────────────────
    s = new()
    _header_bar(s, "Recap 2 · Vision — The Problem",
                "The complaint lives inside the image")
    _kpi_tile(s, Inches(0.5),  Inches(1.3), Inches(3.0), Inches(1.6),
              "Empty-body posts", "3.9%", tint=RGBColor(0xFE,0xF3,0xC7),
              value_color=AMBER)
    _kpi_tile(s, Inches(3.7),  Inches(1.3), Inches(3.0), Inches(1.6),
              "Missed complaints", "text-only", value_color=RED,
              tint=RGBColor(0xFE,0xE2,0xE2))
    _kpi_tile(s, Inches(6.9),  Inches(1.3), Inches(3.0), Inches(1.6),
              "Screenshot share", "~80%", value_color=WALMART_BLUE)
    _kpi_tile(s, Inches(10.1), Inches(1.3), Inches(2.7), Inches(1.6),
              "Impact", "P0", value_color=DARK_BLUE)
    _tx(s, Inches(0.5), Inches(3.15), Inches(12.3), Inches(0.4),
        "Example — a real Walmart-Reddit post", size=Pt(14),
        bold=True, color=DARK_BLUE)
    _table(s, Inches(0.5), Inches(3.55), Inches(12.3), Inches(2.4), [
        ["Field", "Value"],
        ["Title",  "\"Can anyone help me? I need this fixed\""],
        ["Body",   "(empty)"],
        ["Image",  "Screenshot of Walmart-app error message"],
        ["Text-only pipeline",  "scored on title alone — useless"],
        ["Multimodal pipeline", "extracts the error text — understands the complaint"],
    ], first_col_bold=True, highlight_rows={5: RGBColor(0xDC,0xFC,0xE7)})
    _footer(s, page[0], total)

    # ── 10. Recap 2 · Model choice ──────────────────────────────────────────
    s = new()
    _header_bar(s, "Recap 2 · Vision Model — Why Gemma 3 4B")
    _table(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(3.5), [
        ["Model", "DocVQA", "Size", "Ollama", "Verdict"],
        ["gemma3:4b (Google)",  "83",   "3.3 GB", "Yes",  "SELECTED"],
        ["LLaVA-1.5 7B",         "28",   "4.7 GB", "Yes",  "3× worse OCR"],
        ["LLaVA-1.6 8B",         "75",   "5.5 GB", "Yes",  "Larger, slower"],
        ["BLIP-2",               "N/A",  "990 MB", "No",   "Caption only"],
        ["PaliGemma 2 3B",       "81",   "~6 GB",  "No",   "No Ollama runtime"],
    ], first_col_bold=True, highlight_rows={1: RGBColor(0xDC,0xFC,0xE7)})
    _bullets(s, Inches(0.5), Inches(4.9), Inches(12.3), Inches(2.0), [
        "Best DocVQA under 4 GB — reads receipts, screenshots, app screens",
        "Google-maintained (USA) — Walmart vendor-policy compliant",
        "Reuses existing Ollama infrastructure  (localhost:11434)",
        "4–6 s warm latency per image",
    ], size=Pt(13))
    _footer(s, page[0], total)

    # ── 11. Recap 2 · 75% failure rate (villain slide) ──────────────────────
    s = new()
    _header_bar(s, "Recap 2 · The Villain — 75% Failure Rate",
                "First-attempt single-pass captioning was worse than useless")
    _kpi_tile(s, Inches(0.5), Inches(1.3), Inches(3.0), Inches(1.7),
              "Overall failure", "75%",  tint=RGBColor(0xFE,0xE2,0xE2), value_color=RED)
    _kpi_tile(s, Inches(3.7), Inches(1.3), Inches(3.0), Inches(1.7),
              "Hallucination", "50%",   tint=RGBColor(0xFE,0xE2,0xE2), value_color=RED)
    _kpi_tile(s, Inches(6.9), Inches(1.3), Inches(3.0), Inches(1.7),
              "Critical (fake receipts, prices)", "37.5%",
              tint=RGBColor(0xFE,0xE2,0xE2), value_color=RED)
    _kpi_tile(s, Inches(10.1), Inches(1.3), Inches(2.7), Inches(1.7),
              "Correct extraction", "25%", value_color=AMBER)
    _tx(s, Inches(0.5), Inches(3.25), Inches(12.3), Inches(0.4),
        "What hallucination looked like", size=Pt(14), bold=True, color=DARK_BLUE)
    _bullets(s, Inches(0.5), Inches(3.65), Inches(12.3), Inches(2.8), [
        "\"Walmart receipt: $39.99, handwritten 'Damaged Box'\" — actually a product page, no price",
        "\"Order status PENDING\" — no PENDING text anywhere in the image",
        "\"12-pack of 12 fl oz Zero-Sugar Dr Pepper\" — actually a single 42.3 fl oz regular bottle",
        "1 in 2 image posts generated fabricated data → corrupted dashboards.",
    ], size=Pt(13))
    _footer(s, page[0], total)

    # ── 12. Recap 2 · Root cause + papers ───────────────────────────────────
    s = new()
    _header_bar(s, "Recap 2 · Root Cause + Academic Research")
    _tx(s, Inches(0.5), Inches(1.1), Inches(6.0), Inches(0.35),
        "Root cause analysis", size=Pt(14), bold=True, color=DARK_BLUE)
    _bullets(s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(4.5), [
        "768-px fixed resize — small text unreadable",
        "4 B parameters — struggles with multi-element reasoning",
        "No dynamic resolution — one scale for the whole image",
        "Invents plausible details rather than admit uncertainty",
        "No context awareness — screenshot vs physical display",
    ], size=Pt(12))
    _tx(s, Inches(6.8), Inches(1.1), Inches(6.0), Inches(0.35),
        "5 papers reviewed (2023–2025)", size=Pt(14), bold=True, color=DARK_BLUE)
    _table(s, Inches(6.8), Inches(1.5), Inches(6.2), Inches(3.2), [
        ["Paper", "Technique"],
        ["UReader (Tencent)",         "Shape-adaptive cropping"],
        ["TextMonkey (USTC)",         "Shifted window attention"],
        ["DocOwl 1.5 (Alibaba)",      "Structure-aware parsing"],
        ["InternVL2 (Shanghai AI)",   "Tile-based processing"],
        ["Qwen2.5-VL (Alibaba)",      "Dynamic resolution + m-RoPE"],
    ], first_col_bold=True)
    _rect(s, Inches(0.5), Inches(5.6), Inches(12.5), Inches(1.2),
          RGBColor(0xFE,0xF3,0xC7), line=AMBER)
    _tx(s, Inches(0.7), Inches(5.75), Inches(12.1), Inches(0.4),
        "Every model that solves our problem is China-origin — blocked by Walmart vendor policy.",
        size=Pt(13), bold=True, color=AMBER)
    _tx(s, Inches(0.7), Inches(6.2), Inches(12.1), Inches(0.5),
        "Our answer: take the TECHNIQUES from the papers and implement them as CODE on Gemma 3 4B.",
        size=Pt(13), color=DARK_GRAY)
    _footer(s, page[0], total)

    # ── 13. Recap 2 · Multi-pass fix ────────────────────────────────────────
    s = new()
    _header_bar(s, "Recap 2 · The Fix — Multi-Pass Captioning",
                "Same model, same hardware, smarter calling strategy")
    passes = [
        ("PASS 1  ·  STRUCTURE",   "What TYPE of image?\n(screenshot / photo /\nreceipt / app / meme)", PURPLE),
        ("PASS 2  ·  TILE",        "Split into 2–4 crops\n2–4× effective resolution\nno resize", WALMART_BLUE),
        ("PASS 3  ·  EXTRACT",     "Per tile: read ALL\ntext verbatim\n(quotes, prices, buttons)", AMBER),
        ("PASS 4  ·  MERGE",       "Text-only LLM call\n— image is REMOVED —\nmodel cannot invent visuals", GREEN),
    ]
    x = Inches(0.4)
    top = Inches(2.0)
    box_w = Inches(3.0)
    box_h = Inches(2.6)
    gap = Inches(0.15)
    for i, (t, sub, col) in enumerate(passes):
        _rect(s, x, top, box_w, box_h, LIGHT_GRAY, line=col)
        _rect(s, x, top, box_w, Inches(0.5), col)
        _tx(s, x, top + Inches(0.08), box_w, Inches(0.4),
            t, size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _tx(s, x, top + Inches(0.7), box_w, Inches(1.8),
            sub, size=Pt(12), color=DARK_GRAY, align=PP_ALIGN.CENTER)
        if i < 3:
            _arrow(s, x + box_w + Inches(0.005),
                   top + Inches(1.1), Inches(0.14), Inches(0.4))
        x = x + box_w + gap
    _rect(s, Inches(0.5), Inches(5.0), Inches(12.5), Inches(1.6),
          RGBColor(0xDC,0xFC,0xE7), line=GREEN)
    _tx(s, Inches(0.7), Inches(5.15), Inches(12.1), Inches(0.5),
        "Key insight",  size=Pt(13), bold=True, color=GREEN)
    _tx(s, Inches(0.7), Inches(5.55), Inches(12.1), Inches(1.0),
        "By removing the image from the final generation step the model physically cannot invent visual details "
        "— it can only work with text actually extracted in Pass 3.",
        size=Pt(13), color=DARK_GRAY)
    _footer(s, page[0], total)

    # ── 14. Recap 2 · Vision results (payoff) ───────────────────────────────
    s = new()
    _header_bar(s, "Recap 2 · Payoff — Hallucination Eliminated")
    _table(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(2.5), [
        ["Metric", "Before (single-pass)", "After (multi-pass)", "Change"],
        ["Hallucination rate",     "50 % (4/8)",   "0 % (0/8)",   "↓ 100 %"],
        ["Overall failure",        "75 % (6/8)",   "25 % (2/8)",  "↓ 67 %"],
        ["Correct text extraction","25 % (2/8)",   "75 % (6/8)",  "3× better"],
        ["Fabricated claims",      "8 total",      "0",           "eliminated"],
        ["Latency / image",        "~5 s",         "~15 s",       "3× (acceptable)"],
    ], first_col_bold=True, highlight_rows={1: RGBColor(0xDC,0xFC,0xE7),
                                             3: RGBColor(0xDC,0xFC,0xE7)})
    _tx(s, Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.4),
        "Scaled validation — 25 images", size=Pt(14), bold=True, color=DARK_BLUE)
    _kpi_tile(s, Inches(0.5), Inches(4.45), Inches(3.9), Inches(1.9),
              "PASS", "22 / 25 (88%)", value_color=GREEN,
              tint=RGBColor(0xDC,0xFC,0xE7))
    _kpi_tile(s, Inches(4.7), Inches(4.45), Inches(3.9), Inches(1.9),
              "PARTIAL", "3 / 25 (12%)", value_color=AMBER,
              tint=RGBColor(0xFE,0xF3,0xC7))
    _kpi_tile(s, Inches(8.9), Inches(4.45), Inches(3.9), Inches(1.9),
              "FAIL", "0 / 25 (0%)", value_color=GREEN,
              tint=RGBColor(0xDC,0xFC,0xE7))
    _footer(s, page[0], total)

    # ── 15. Recap 3 · Trust score ───────────────────────────────────────────
    s = new()
    _header_bar(s, "Recap 3 · Trust Score — Weighted Formula")
    _rect(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.0),
          LIGHT_BLUE, line=WALMART_BLUE)
    _tx(s, Inches(0.7), Inches(1.4), Inches(12.0), Inches(0.6),
        "trust_score  =  0.4 × metadata  +  0.3 × dedup  +  0.3 × llm_credibility",
        size=Pt(18), bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    _tx(s, Inches(0.5), Inches(2.4), Inches(12.3), Inches(0.4),
        "Component weights inside metadata score",
        size=Pt(13), bold=True, color=DARK_BLUE)
    _table(s, Inches(0.5), Inches(2.85), Inches(12.3), Inches(3.4), [
        ["Signal",           "Formula",                                                          "Weight"],
        ["Base floor",       "constant",                                                         "0.15"],
        ["Account age",      "min(account_age_days / 365, 1.0)",                                 "0.20"],
        ["Karma",            "min(total_karma / 5000, 1.0)",                                     "0.20"],
        ["Content length",   "min((len(title) + len(body)) / 200, 1.0)",                         "0.30"],
        ["Engagement",       "min(max(reddit_score, 0) / 20, 1.0)",                              "0.15"],
    ], first_col_bold=True)
    _tx(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
        "LLM credibility is invoked ONLY in the ambiguous zone  (0.3 < metadata < 0.8)  → cost control",
        size=Pt(11), color=MED_GRAY)
    _footer(s, page[0], total)

    # ── 16. Recap 3 · Confidence + priority ─────────────────────────────────
    s = new()
    _header_bar(s, "Recap 3 · Confidence & Priority Classification")
    _tx(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.4),
        "Confidence = softmax probability of the predicted class",
        size=Pt(14), bold=True, color=DARK_BLUE)
    _table(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(2.0), [
        ["Threshold",             "Value",                              "Source"],
        ["Analysis confidence",   "≥ 0.70",                             "config/models.yaml"],
        ["Notification P1",       "trust ≥ 0.70  AND  conf ≥ 0.80",     "dispatcher.py"],
        ["Notification P2",       "trust ≥ 0.50  AND  conf ≥ 0.60",     "dispatcher.py"],
    ], first_col_bold=True)
    # P1 / P2 cards
    _rect(s, Inches(0.5), Inches(4.0), Inches(6.0), Inches(2.6),
          RGBColor(0xFE,0xE2,0xE2), line=RED)
    _tx(s, Inches(0.7), Inches(4.2), Inches(5.6), Inches(0.5),
        "P1  ·  Immediate action", size=Pt(18), bold=True, color=RED)
    _tx(s, Inches(0.7), Inches(4.75), Inches(5.6), Inches(1.6),
        "trust ≥ 0.70  ∧  confidence ≥ 0.80\nHigh-signal negative posts\nrouted to on-call groups",
        size=Pt(13), color=DARK_GRAY)
    _rect(s, Inches(6.8), Inches(4.0), Inches(6.0), Inches(2.6),
          RGBColor(0xFE,0xF3,0xC7), line=AMBER)
    _tx(s, Inches(7.0), Inches(4.2), Inches(5.6), Inches(0.5),
        "P2  ·  Review-worthy", size=Pt(18), bold=True, color=AMBER)
    _tx(s, Inches(7.0), Inches(4.75), Inches(5.6), Inches(1.6),
        "trust ≥ 0.50  ∧  confidence ≥ 0.60\nQueue for analyst review\nlower urgency",
        size=Pt(13), color=DARK_GRAY)
    _footer(s, page[0], total)

    # ── 17. Recap 4 · Why ModernBERT ────────────────────────────────────────
    s = new()
    _header_bar(s, "Recap 4 · Why ModernBERT — not RoBERTa",
                "Long-context, domain-fine-tuned encoders beat Twitter baselines on Reddit complaints")
    _table(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(4.2), [
        ["Criterion",                "RoBERTa (cardiffnlp/twitter)",  "ModernBERT (answerdotai)",  "Winner"],
        ["Context length",           "512 tokens",                     "8192 tokens  (16×)",         "ModernBERT"],
        ["Training corpus",          "Twitter (short, casual)",        "Web + code + long docs",     "ModernBERT"],
        ["Long-post F1 (≥512 tok)",  "0.2778",                         "1.0000",                     "ModernBERT"],
        ["Overall Macro F1",         "0.6272",                         "0.7642",                     "ModernBERT"],
        ["Latency (MPS, ms/post)",   "6.5 ms",                         "11.9 ms",                    "RoBERTa"],
        ["Off-the-shelf quality",    "Strong baseline",                "Weak without fine-tuning",   "RoBERTa"],
    ], first_col_bold=True, highlight_rows={3: RGBColor(0xDC,0xFC,0xE7),
                                             4: RGBColor(0xDC,0xFC,0xE7)})
    _rect(s, Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.2),
          LIGHT_BLUE, line=WALMART_BLUE)
    _tx(s, Inches(0.7), Inches(5.75), Inches(12.0), Inches(0.4),
        "Decision", size=Pt(13), bold=True, color=WALMART_BLUE)
    _tx(s, Inches(0.7), Inches(6.15), Inches(12.0), Inches(0.5),
        "Accept +5 ms latency for +0.137 Macro F1.  RoBERTa is retained as fallback in config/models.yaml if the local checkpoint is missing.",
        size=Pt(12), color=DARK_GRAY)
    _footer(s, page[0], total)

    # ── 18. Recap 4 · 3-stage curriculum ────────────────────────────────────
    s = new()
    _header_bar(s, "Recap 4 · 3-Stage Curriculum Training")
    stages = [
        ("STAGE 1", "TweetEval-sentiment\n45 K tweets", "2 epochs\nPolarity grounding",  AMBER),
        ("STAGE 2", "GoEmotions-3class\n54 K Reddit",   "2 epochs\nReddit register",     WALMART_BLUE),
        ("STAGE 3", "Walmart-200\n5-fold CV",           "up to 15 epochs\npatience=3\nDomain specialization", GREEN),
    ]
    x = Inches(0.6)
    top = Inches(1.5)
    box_w = Inches(4.0)
    box_h = Inches(2.6)
    gap = Inches(0.15)
    for i, (n, dataset, purpose, col) in enumerate(stages):
        _rect(s, x, top, box_w, box_h, LIGHT_GRAY, line=col)
        _rect(s, x, top, box_w, Inches(0.5), col)
        _tx(s, x, top + Inches(0.08), box_w, Inches(0.4),
            n, size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _tx(s, x, top + Inches(0.65), box_w, Inches(0.8),
            dataset, size=Pt(13), bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
        _tx(s, x, top + Inches(1.5), box_w, Inches(1.0),
            purpose, size=Pt(12), color=DARK_GRAY, align=PP_ALIGN.CENTER)
        if i < 2:
            _arrow(s, x + box_w + Inches(0.005),
                   top + Inches(1.1), Inches(0.14), Inches(0.4))
        x = x + box_w + gap
    _tx(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.4),
        "Training configuration", size=Pt(13), bold=True, color=DARK_BLUE)
    _bullets(s, Inches(0.5), Inches(4.9), Inches(12.3), Inches(2.0), [
        "max_length = 1024 tokens  (key lever for long-context advantage)",
        "Effective batch size 32  (per-device BS=8 × grad-accum=4)",
        "Class weights  neg=0.52  neu=1.03  pos=8.33  (inverse frequency)",
        "Minority oversampling to ~100/class per training fold",
        "Early stopping on eval_macro_f1 with patience=3",
        "Hardware — Apple M-series (MPS backend)",
    ], size=Pt(12))
    _footer(s, page[0], total)

    # ── 19. Recap 4 · Final results (payoff) ────────────────────────────────
    s = new()
    _header_bar(s, "Recap 4 · Payoff — ModernBERT vs RoBERTa",
                "5-fold out-of-fold cross-validation on 200 real Walmart-Reddit posts")
    _kpi_tile(s, Inches(0.5), Inches(1.2), Inches(4.0), Inches(1.8),
              "Overall Macro F1", "0.6272 → 0.7642", value_color=GREEN,
              tint=RGBColor(0xDC,0xFC,0xE7))
    _kpi_tile(s, Inches(4.7), Inches(1.2), Inches(4.0), Inches(1.8),
              "Long-post F1  (≥ 512 tok)", "0.28 → 1.00",
              value_color=GREEN, tint=RGBColor(0xDC,0xFC,0xE7))
    _kpi_tile(s, Inches(8.9), Inches(1.2), Inches(3.9), Inches(1.8),
              "Improvement", "+22 %", value_color=WALMART_BLUE)
    _table(s, Inches(0.5), Inches(3.2), Inches(12.3), Inches(3.6), [
        ["Metric",                        "RoBERTa",  "ModernBERT v2",  "Δ"],
        ["Macro F1 (overall)",             "0.6272",  "0.7642",         "+0.137"],
        ["F1 negative",                    "0.7967",  "0.8779",         "+0.081"],
        ["F1 neutral",                     "0.6087",  "0.7480",         "+0.139"],
        ["F1 positive",                    "0.4762",  "0.6667",         "+0.190"],
        ["Long-post F1 (≥ 512 tokens)",    "0.2778",  "1.0000",         "+0.722"],
        ["Short-post F1 (n = 193)",        "0.6360",  "0.7619",         "+0.126"],
        ["Latency (ms/post, MPS)",         "6.5",     "11.9",           "+5.4 ms"],
    ], first_col_bold=True, highlight_rows={5: RGBColor(0xDC,0xFC,0xE7)})
    _footer(s, page[0], total)

    # ── 20. Recap 4 · Training Evidence (NEW) ───────────────────────────────
    s = new()
    _header_bar(s, "Recap 4 · Training Evidence — Artifacts on Disk",
                "Proof ModernBERT was actually fine-tuned in-house (not a downloaded checkpoint)")
    _tx(s, Inches(0.5), Inches(1.1), Inches(6.2), Inches(0.35),
        "Training pipeline (scripts → outputs)", size=Pt(13), bold=True, color=DARK_BLUE)
    _table(s, Inches(0.5), Inches(1.5), Inches(6.3), Inches(4.3), [
        ["Stage",           "Script"],
        ["Data collection", "fetch_real_benchmark.py"],
        ["Human labeling",  "label_benchmark.py"],
        ["Curriculum train","train_modernbert_sentiment.py"],
        ["Honest eval",     "eval_sentiment_models.py"],
        ["Thesis chapter",  "docs/MODEL_COMPARISON.md"],
    ], first_col_bold=True, body_font=Pt(10))
    _tx(s, Inches(7.0), Inches(1.1), Inches(5.8), Inches(0.35),
        "Checkpoints produced (on disk today)", size=Pt(13), bold=True, color=DARK_BLUE)
    _rect(s, Inches(7.0), Inches(1.5), Inches(5.9), Inches(3.3), DARK_BLUE)
    _tx(s, Inches(7.15), Inches(1.6), Inches(5.7), Inches(3.1), [
        "models/modernbert_walmart/",
        "├─ stage1_tweeteval/     macro F1 = 0.7267",
        "├─ stage2_goemotions/    macro F1 = 0.7028",
        "├─ stage3_walmart/       5-fold CV artefacts",
        "├─ final/                production checkpoint",
        "├─ final_max512/         v1 ablation",
        "└─ eval_results.json     aggregated CV metrics",
    ], size=Pt(11), color=WHITE)
    # Reproduction command box
    _rect(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.0),
          DARK_GRAY)
    _tx(s, Inches(0.7), Inches(6.05), Inches(12.0), Inches(0.35),
        "Reproduction (offline)", size=Pt(10), bold=True, color=YELLOW)
    _tx(s, Inches(0.7), Inches(6.4), Inches(12.0), Inches(0.55),
        "HF_HUB_OFFLINE=1 TRANSFORMERS_OFFLINE=1 HF_DATASETS_OFFLINE=1  python scripts/train_modernbert_sentiment.py  --stages 1,2,3 --folds 5 --max-length 1024 --batch-size 8",
        size=Pt(10), color=WHITE)
    # Smoke test note
    _tx(s, Inches(0.5), Inches(5.9), Inches(6.3), Inches(0.4),
        "Integration smoke test — 5/5 correct on first 5 real benchmark posts",
        size=Pt(11), bold=True, color=GREEN)
    _footer(s, page[0], total)

    # ── 21. Recap 5 · Dashboard pages ───────────────────────────────────────
    s = new()
    _header_bar(s, "Recap 5 · Dashboard Pages")
    _table(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(4.6), [
        ["Page",              "Priority", "Purpose",                                     "Key metric"],
        ["Brand Health",      "P0",       "At-a-glance KPIs & trends",                   "Overall sentiment score"],
        ["Post Explorer",     "P1",       "Search / filter all posts",                   "Volume + sentiment distribution"],
        ["Review & Validate", "P0",       "Correct labels + draft replies",              "Accuracy improvement"],
        ["Post Lifecycle",    "P0",       "Kanban (triaged → resolved)",                 "Resolution rate"],
        ["Insights",          "P1",       "AI-generated competitor analysis",            "Issue rankings"],
        ["Pipeline Control",  "P1",       "Monitor & trigger runs",                      "Jobs, cursors, health"],
        ["Notifications",     "P1",       "Group-based alert routing config",            "Delivery log"],
    ], first_col_bold=True, body_font=Pt(11))
    _tx(s, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.4),
        "WebSocket push for real-time alerts — no polling needed",
        size=Pt(11), bold=True, color=WALMART_BLUE)
    _footer(s, page[0], total)

    # ── 22. Recap 6 · Notification system ───────────────────────────────────
    s = new()
    _header_bar(s, "Recap 6 · Notification System",
                "Group-based routing for P1 / P2 posts")
    steps = [
        ("Pipeline\nanalyses post",  WALMART_BLUE),
        ("classify_priority\n(trust, conf)", PURPLE),
        ("Match subreddit\n+ priority filter", AMBER),
        ("Send email +\nSlack per group", GREEN),
        ("Log to\nnotification_log", DARK_BLUE),
    ]
    x = Inches(0.4)
    top = Inches(1.6)
    box_w = Inches(2.4)
    box_h = Inches(1.6)
    gap = Inches(0.15)
    for i, (t, col) in enumerate(steps):
        _rect(s, x, top, box_w, box_h, LIGHT_GRAY, line=col)
        _rect(s, x, top, box_w, Inches(0.15), col)
        _tx(s, x, top + Inches(0.35), box_w, Inches(1.2),
            t, size=Pt(13), bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE)
        if i < 4:
            _arrow(s, x + box_w + Inches(0.005),
                   top + Inches(0.65), Inches(0.14), Inches(0.35))
        x = x + box_w + gap
    _tx(s, Inches(0.5), Inches(3.55), Inches(12.3), Inches(0.4),
        "Configuration page (/notifications)  ·  8 REST endpoints",
        size=Pt(14), bold=True, color=DARK_BLUE)
    _bullets(s, Inches(0.5), Inches(4.0), Inches(12.3), Inches(2.8), [
        "Create groups — name, subreddits, email DL, Slack channel",
        "Quick-add subreddits by category (Walmart core, Spark, Pharmacy, International, Sam's Club, Competitors)",
        "Priority filter — choose P1, P2, or both per group",
        "Enable / disable toggle + Test (dry-run) simulator",
        "Delivery log — audit trail of all sent notifications",
    ], size=Pt(12))
    _footer(s, page[0], total)

    # ── 23. Part B divider ──────────────────────────────────────────────────
    s = new()
    _set_bg(s, LIGHT_BLUE)
    _rect(s, Inches(0), Inches(3.0), SLIDE_W, Inches(1.5), WALMART_BLUE)
    _tx(s, Inches(0.5), Inches(3.15), Inches(12.3), Inches(0.7),
        "PART B", size=Pt(30), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _tx(s, Inches(0.5), Inches(3.65), Inches(12.3), Inches(0.7),
        "Post Mid-Semester Work", size=Pt(40), bold=True, color=WHITE,
        align=PP_ALIGN.CENTER)
    _tx(s, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.5),
        "5 new capabilities · analyst workflows · lifecycle · competitor insights",
        size=Pt(14), color=DARK_GRAY, align=PP_ALIGN.CENTER)
    _footer(s, page[0], total)

    # ── 24. Post-midsem 1 · Review & Validate ───────────────────────────────
    s = new()
    _header_bar(s, "Post-Midsem 1 · Review & Validate",
                "Human-in-the-loop correction and reply generation")
    _tx(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.4),
        "Workflow", size=Pt(14), bold=True, color=DARK_BLUE)
    _bullets(s, Inches(0.5), Inches(1.55), Inches(12.3), Inches(1.6), [
        "Queue shows posts sorted by priority (P1 first)",
        "Analyst reviews sentiment + aspects (correct if wrong)",
        "Click \"Generate Drafts\" → two reply options generated",
        "Analyst picks A or B, edits, and posts to Reddit",
    ], size=Pt(13))
    # Dual-draft cards
    _rect(s, Inches(0.5), Inches(3.4), Inches(6.0), Inches(2.8),
          LIGHT_BLUE, line=WALMART_BLUE)
    _tx(s, Inches(0.7), Inches(3.55), Inches(5.7), Inches(0.4),
        "Draft A · Smart Composer", size=Pt(14), bold=True, color=WALMART_BLUE)
    _tx(s, Inches(0.7), Inches(3.95), Inches(5.7), Inches(2.0),
        "Keyword extraction  +  curated phrase pools\nDeterministic, always available\nGood default tone for retail complaints",
        size=Pt(12), color=DARK_GRAY)
    _rect(s, Inches(6.8), Inches(3.4), Inches(6.0), Inches(2.8),
          RGBColor(0xED, 0xE9, 0xFE), line=PURPLE)
    _tx(s, Inches(7.0), Inches(3.55), Inches(5.7), Inches(0.4),
        "Draft B · FLAN-T5-base", size=Pt(14), bold=True, color=PURPLE)
    _tx(s, Inches(7.0), Inches(3.95), Inches(5.7), Inches(2.0),
        "Multi-temperature sampling  +  scorer\nBest-of-N selection\nHigher variety, less templated",
        size=Pt(12), color=DARK_GRAY)
    _tx(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.4),
        "Learning loop — posted replies saved to feedback table → future few-shot examples",
        size=Pt(11), bold=True, color=GREEN)
    _footer(s, page[0], total)

    # ── 25. Post-midsem 2 · Post Explorer ───────────────────────────────────
    s = new()
    _header_bar(s, "Post-Midsem 2 · Post Explorer",
                "Search, filter, and deep-dive into analyzed posts")
    _tx(s, Inches(0.5), Inches(1.15), Inches(6.0), Inches(0.4),
        "Filters available", size=Pt(14), bold=True, color=DARK_BLUE)
    _bullets(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(4.5), [
        "Sentiment · negative / neutral / positive",
        "Confidence threshold slider",
        "Trust-score threshold slider",
        "Subreddit — multi-select from 25 tracked",
        "Aspect — delivery, product, returns, support, price, app",
        "Date range — today / week / month / custom",
        "Full-text search across titles & bodies",
    ], size=Pt(13))
    _tx(s, Inches(6.8), Inches(1.15), Inches(6.0), Inches(0.4),
        "Per-post card shows", size=Pt(14), bold=True, color=DARK_BLUE)
    _bullets(s, Inches(6.8), Inches(1.55), Inches(6.0), Inches(4.5), [
        "Title + body excerpt",
        "Sentiment badge (color-coded) + confidence %",
        "Trust-score indicator",
        "Aspect tags (multi-aspect supported)",
        "Subreddit + post time",
        "Actions — Review · Add to Lifecycle · View Details",
    ], size=Pt(13))
    _footer(s, page[0], total)

    # ── 26. Post-midsem 3 · Lifecycle Kanban ────────────────────────────────
    s = new()
    _header_bar(s, "Post-Midsem 3 · Post Lifecycle — Kanban",
                "Track a complaint from New → Resolved with SLA visibility")
    states = [
        ("TRIAGED",       "New P1/P2\nposts land here",           AMBER),
        ("ACKNOWLEDGED",  "Assigned\nto analyst",                 WALMART_BLUE),
        ("IN PROGRESS",   "Reply being\ndrafted",                 PURPLE),
        ("RESOLVED",      "Reply posted\nSLA closed",             GREEN),
    ]
    x = Inches(0.5)
    top = Inches(1.7)
    box_w = Inches(3.0)
    box_h = Inches(2.4)
    gap = Inches(0.1)
    for i, (t, sub, col) in enumerate(states):
        _rect(s, x, top, box_w, box_h, LIGHT_GRAY, line=col)
        _rect(s, x, top, box_w, Inches(0.5), col)
        _tx(s, x, top + Inches(0.1), box_w, Inches(0.4),
            t, size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _tx(s, x, top + Inches(0.9), box_w, Inches(1.3),
            sub, size=Pt(12), color=DARK_GRAY, align=PP_ALIGN.CENTER)
        if i < 3:
            _arrow(s, x + box_w + Inches(0.005),
                   top + Inches(1.0), Inches(0.1), Inches(0.4))
        x = x + box_w + gap
    _tx(s, Inches(0.5), Inches(4.4), Inches(12.3), Inches(0.4),
        "Resolve modal — 2-step flow", size=Pt(14), bold=True, color=DARK_BLUE)
    _bullets(s, Inches(0.5), Inches(4.85), Inches(12.3), Inches(2.0), [
        "Step 1  —  Save action note + optional LLM-drafted reply",
        "         a) Save reply & open Reddit  (copies to clipboard)",
        "         b) OR  Resolve (no reply needed)  — close without posting",
        "Step 2  —  Paste reply on Reddit → return to dashboard → \"Mark Resolved\"",
        "All state transitions are logged with timestamps for SLA tracking",
    ], size=Pt(12))
    _footer(s, page[0], total)

    # ── 27. Post-midsem 4 · Insights & competitors ──────────────────────────
    s = new()
    _header_bar(s, "Post-Midsem 4 · Insights & Competitor Analysis",
                "AI-generated strategic intelligence from raw data")
    features = [
        ("Issue Rankings",     "Top negative issues ranked by\nvolume × severity × recency\ngrouped by aspect",  WALMART_BLUE),
        ("Competitor Pulse",   "Walmart vs Costco, Target,\nAmazon on shared aspects\ncross-mentioned posts",     PURPLE),
        ("LLM Summarization",  "Natural-language weekly summaries\nsuggested action items\nemerging-topic detection", GREEN),
        ("Aspect Drilldown",   "6-category taxonomy\nper-aspect sentiment trend + volume\nrepresentative posts",   AMBER),
    ]
    # 2x2 grid
    positions = [(0.5, 1.3), (6.9, 1.3), (0.5, 4.0), (6.9, 4.0)]
    for (t, sub, col), (lx, ly) in zip(features, positions):
        _rect(s, Inches(lx), Inches(ly), Inches(6.0), Inches(2.5),
              LIGHT_GRAY, line=col)
        _rect(s, Inches(lx), Inches(ly), Inches(6.0), Inches(0.5), col)
        _tx(s, Inches(lx), Inches(ly + 0.1), Inches(6.0), Inches(0.4),
            t, size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _tx(s, Inches(lx + 0.2), Inches(ly + 0.7), Inches(5.6), Inches(1.7),
            sub, size=Pt(12), color=DARK_GRAY, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE)
    _footer(s, page[0], total)

    # ── 28. Live Demo placeholder ───────────────────────────────────────────
    s = new()
    _set_bg(s, DARK_BLUE)
    _tx(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(1.0),
        "LIVE DEMO", size=Pt(60), bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
    _tx(s, Inches(0.5), Inches(3.4), Inches(12.3), Inches(0.6),
        "Run one pipeline cycle · watch an alert fire in the dashboard",
        size=Pt(20), color=WHITE, align=PP_ALIGN.CENTER)
    _rect(s, Inches(1.5), Inches(4.5), Inches(10.3), Inches(2.0),
          LIGHT_BLUE, line=YELLOW)
    _tx(s, Inches(1.7), Inches(4.7), Inches(10.0), Inches(0.4),
        "Demo checklist", size=Pt(14), bold=True, color=DARK_BLUE)
    _bullets(s, Inches(1.7), Inches(5.15), Inches(10.0), Inches(1.5), [
        "Trigger \"Run Now\" from Pipeline page  ·  show scheduler ticker",
        "Open Brand Health → new tile counts update after run",
        "Open Alerts → new P1 alert lands via WebSocket (no refresh)",
        "Open the alerted post in Review & Validate → generate + edit reply",
    ], size=Pt(13), color=DARK_GRAY)
    _footer(s, page[0], total)

    # ── 29. Post-midsem 5 · Consolidated results ────────────────────────────
    s = new()
    _header_bar(s, "Post-Midsem 5 · Key Achievements",
                "Consolidated numbers across all subsystems")
    _table(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.2), [
        ["Area",           "Achievement",                                      "Evidence"],
        ["ModernBERT",     "Macro F1  0.6272 → 0.7642  (+22 %)",               "5-fold OOF CV, n=200"],
        ["Long-context",   "Long-post F1  0.28 → 1.00  (+722 %)",              "7 posts ≥ 512 tokens"],
        ["Vision",         "Hallucination  50 % → 0 %",                        "8 + 25 image validation"],
        ["Vision",         "Text extraction  25 % → 75 %",                     "Multi-pass pipeline"],
        ["Pipeline",       "25 subreddits, every 6 h (automated)",             "Arctic Shift API"],
        ["Dashboard",      "7 pages, real-time WebSocket alerts",              "React + FastAPI"],
        ["Notifications",  "Group-based P1 / P2 routing",                      "Email + Slack channels"],
        ["Lifecycle",      "Full Kanban with 2-step resolve",                  "Triage → Resolved"],
    ], first_col_bold=True, body_font=Pt(11),
       highlight_rows={1: RGBColor(0xDC,0xFC,0xE7),
                       2: RGBColor(0xDC,0xFC,0xE7),
                       3: RGBColor(0xDC,0xFC,0xE7)})
    _footer(s, page[0], total)

    # ── 30. Technical stack ─────────────────────────────────────────────────
    s = new()
    _header_bar(s, "Technical Stack Summary")
    _table(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.2), [
        ["Layer",         "Technology",                                          "Key decision"],
        ["Backend",       "Python 3.13 · FastAPI · SQLite",                      "Free, local-first, modular"],
        ["Frontend",      "React 18 · TypeScript · Vite · Tailwind",             "Modern SPA, responsive"],
        ["Sentiment",     "ModernBERT (fine-tuned, 1024 tokens)",                "Domain-specialized, offline"],
        ["Aspects",       "BART-MNLI  (zero-shot)",                              "No training needed"],
        ["Vision",        "Gemma 3 4B via Ollama (multi-pass)",                  "Policy compliant, 0 % hallucination"],
        ["Reply Gen",     "FLAN-T5  +  Smart Composer  (dual draft)",            "Learning loop via feedback"],
        ["Trust",         "metadata + dedup + LLM (weighted)",                   "Flag, don't drop"],
        ["Scheduling",    "asyncio lifespan (6 h)  +  manual",                   "Cursor-based incremental"],
        ["Observability", "structlog  +  cost ledger (JSONL)",                   "Per-call LLM cost tracking"],
    ], first_col_bold=True, body_font=Pt(10))
    _footer(s, page[0], total)

    # ── 31. Future work ─────────────────────────────────────────────────────
    s = new()
    _header_bar(s, "Future Work — Semester 5 Roadmap")
    _bullets(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.5), [
        "Azure Cosmos DB migration — production-grade storage",
        "Twitter / X integration as second data source",
        "3-seed ensemble for ModernBERT  (+0.01–0.03 F1, tighter variance)",
        "Gemma 3 12B upgrade for remaining edge-case images",
        "Spanish language support (bilingual retail communities)",
        "Azure AD authentication for multi-user access",
        "Automated retraining pipeline — feedback loop → model updates",
        "Slack bot integration for inline notification responses",
        "Blind 25-post recheck for ModernBERT defensibility",
    ], size=Pt(14))
    _footer(s, page[0], total)

    # ── 32. Thank you ───────────────────────────────────────────────────────
    s = new()
    _set_bg(s, DARK_BLUE)
    _tx(s, Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.2),
        "Thank You", size=Pt(72), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _tx(s, Inches(0.5), Inches(3.9), Inches(12.3), Inches(0.5),
        "Retail Sentiment Intelligence  ·  Post Mid-Semester Progress",
        size=Pt(20), color=YELLOW, align=PP_ALIGN.CENTER)
    _tx(s, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.5),
        "Questions ?", size=Pt(28), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _tx(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.4),
        "Vishal Singh  ·  2020AA05641  ·  BITS Pilani (WILP)  ·  Walmart Global Tech, Bengaluru",
        size=Pt(12), color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

    # ── Save ────────────────────────────────────────────────────────────────
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"[OK] wrote {OUT}  ({page[0]} slides)")


if __name__ == "__main__":
    build()
