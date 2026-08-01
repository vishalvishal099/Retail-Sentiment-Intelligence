"""
Generate the FINAL post-midsem dissertation presentation.

Focus: only post-midsem work — no mid-sem recap. Content pulled from
the final dissertation report (docs/Sem_4/final/latex/) and the UI
screenshots captured for the report (docs/Sem_4/final/figures/ui/).

Output: docs/Sem_4/FINAL_PRESENTATION_VishalSingh_2020AA05641.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = Path(__file__).parent
FIG_DIR = HERE / "figures"
UI_DIR = HERE / "final" / "figures" / "ui"
OUT = HERE / "FINAL_PRESENTATION_VishalSingh_2020AA05641.pptx"

# Palette — matches the mid-sem deck for continuity.
WALMART_BLUE = RGBColor(0x00, 0x71, 0xDC)
DARK_BLUE    = RGBColor(0x04, 0x1E, 0x42)
LIGHT_BLUE   = RGBColor(0xE8, 0xF4, 0xFD)
YELLOW       = RGBColor(0xFF, 0xC2, 0x20)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY    = RGBColor(0x33, 0x33, 0x33)
MED_GRAY     = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY   = RGBColor(0xF5, 0xF5, 0xF5)
GREEN        = RGBColor(0x10, 0xB9, 0x81)
RED          = RGBColor(0xEF, 0x44, 0x44)
PURPLE       = RGBColor(0x7C, 0x3A, 0xED)
AMBER        = RGBColor(0xD9, 0x77, 0x06)
GREEN_TINT   = RGBColor(0xDC, 0xFC, 0xE7)
AMBER_TINT   = RGBColor(0xFE, 0xF3, 0xC7)
RED_TINT     = RGBColor(0xFE, 0xE2, 0xE2)
PURPLE_TINT  = RGBColor(0xED, 0xE9, 0xFE)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ─── helpers ─────────────────────────────────────────────────────────────────
def _set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _tx(slide, left, top, width, height, text, *,
        size=Pt(14), bold=False, color=DARK_GRAY,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        for r in p.runs:
            r.font.size = size
            r.font.bold = bold
            r.font.color.rgb = color
    return box


def _rect(slide, left, top, width, height, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def _header_bar(slide, title, subtitle=None):
    _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.85), DARK_BLUE)
    _tx(slide, Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.5),
        title, size=Pt(26), bold=True, color=WHITE)
    if subtitle:
        _tx(slide, Inches(0.4), Inches(0.5), Inches(12.5), Inches(0.35),
            subtitle, size=Pt(12), color=LIGHT_BLUE)
    _rect(slide, Inches(0), Inches(0.85), SLIDE_W, Inches(0.05), YELLOW)


def _footer(slide, page_no, total):
    _tx(slide, Inches(0.4), Inches(7.15), Inches(9), Inches(0.3),
        "Retail Sentiment Intelligence  ·  BITS ZG628T Dissertation  ·  Vishal Singh  ·  Post Mid-Semester Presentation",
        size=Pt(9), color=MED_GRAY)
    _tx(slide, Inches(12.0), Inches(7.15), Inches(1.2), Inches(0.3),
        f"{page_no} / {total}", size=Pt(9), color=MED_GRAY, align=PP_ALIGN.RIGHT)


def _bullets(slide, left, top, width, height, items, *,
             size=Pt(14), color=DARK_GRAY, bold=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {item}"
        p.space_after = Pt(4)
        for r in p.runs:
            r.font.size = size
            r.font.color.rgb = color
            r.font.bold = bold
    return box


def _table(slide, left, top, width, height, data, *,
           header_fill=DARK_BLUE, header_color=WHITE,
           body_font=Pt(11), first_col_bold=False,
           highlight_rows=None):
    highlight_rows = highlight_rows or {}
    rows = len(data)
    cols = len(data[0])
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = tbl_shape.table
    for r_idx, row in enumerate(data):
        for c_idx, val in enumerate(row):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = str(val)
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
            elif r_idx in highlight_rows:
                cell.fill.solid()
                cell.fill.fore_color.rgb = highlight_rows[r_idx]
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r_idx % 2 else LIGHT_GRAY
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = body_font
                    if r_idx == 0:
                        run.font.color.rgb = header_color
                        run.font.bold = True
                    else:
                        run.font.color.rgb = DARK_GRAY
                        if first_col_bold and c_idx == 0:
                            run.font.bold = True
    return tbl


def _kpi_tile(slide, left, top, width, height, label, value, *,
              tint=LIGHT_BLUE, value_color=WALMART_BLUE, value_size=Pt(28)):
    _rect(slide, left, top, width, height, tint, line=WALMART_BLUE)
    _tx(slide, left, top + Inches(0.15), width, Inches(0.35),
        label, size=Pt(11), bold=True, color=MED_GRAY, align=PP_ALIGN.CENTER)
    _tx(slide, left, top + Inches(0.55), width, height - Inches(0.7),
        value, size=value_size, bold=True, color=value_color,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def _arrow(slide, x, y, w=Inches(0.35), h=Inches(0.35), color=YELLOW):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    a.fill.solid()
    a.fill.fore_color.rgb = color
    a.line.fill.background()


def _stage_card(slide, x, y, w, h, label, title, body, color):
    _rect(slide, x, y, w, h, LIGHT_GRAY, line=color)
    _rect(slide, x, y, w, Inches(0.45), color)
    _tx(slide, x, y + Inches(0.05), w, Inches(0.35),
        label, size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _tx(slide, x, y + Inches(0.55), w, Inches(0.4),
        title, size=Pt(13), bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    _tx(slide, x, y + Inches(1.0), w, h - Inches(1.05),
        body, size=Pt(11), color=DARK_GRAY, align=PP_ALIGN.CENTER)


def _image(slide, path: Path, left, top, width=None, height=None):
    if not path.exists():
        _rect(slide, left, top, width or Inches(6), height or Inches(3),
              LIGHT_GRAY, line=MED_GRAY)
        _tx(slide, left, top, width or Inches(6), height or Inches(3),
            f"[missing image: {path.name}]", size=Pt(11), color=MED_GRAY,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        return
    kw = {}
    if width is not None:
        kw["width"] = width
    if height is not None:
        kw["height"] = height
    slide.shapes.add_picture(str(path), left, top, **kw)


def _link(slide, left, top, width, height, url, label=None, *,
          size=Pt(10), bold=False, color=WALMART_BLUE, align=PP_ALIGN.LEFT):
    """Add a clickable hyperlink text box."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = label or url
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.underline = True
    run.hyperlink.address = url
    return box


# ─── deck ────────────────────────────────────────────────────────────────────
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    TOTAL = 22
    page = [0]

    def new():
        page[0] += 1
        return prs.slides.add_slide(prs.slide_layouts[6])

    # 1 · Title ──────────────────────────────────────────────────────────────
    s = new()
    _set_bg(s, DARK_BLUE)
    _rect(s, Inches(0), Inches(2.3), Inches(0.4), Inches(2.7), YELLOW)
    _tx(s, Inches(0.9), Inches(0.85), Inches(11.5), Inches(0.9),
        "Retail Sentiment Intelligence", size=Pt(38), bold=True, color=WHITE)
    _tx(s, Inches(0.9), Inches(1.75), Inches(11.5), Inches(1.4),
        "Real-Time Social Media Mining and Trust-Aware Sentiment Analysis\n"
        "Using Large Language Models for Retail Feedback Optimization",
        size=Pt(18), color=YELLOW)
    _tx(s, Inches(0.9), Inches(3.15), Inches(11.5), Inches(0.4),
        "BITS ZG628T  ·  Dissertation  ·  Post Mid-Semester Presentation",
        size=Pt(14), bold=True, color=WHITE)
    # People — Faculty Mentor → Supervisor → Candidate (candidate at the end)
    _tx(s, Inches(0.9), Inches(4.15), Inches(11.5), Inches(0.4),
        "Faculty Mentor:   Ms. Pradnya Kashikar",
        size=Pt(15), bold=True, color=WHITE)
    _tx(s, Inches(0.9), Inches(4.6), Inches(11.5), Inches(0.4),
        "BITS Pilani (WILP)",
        size=Pt(12), color=LIGHT_BLUE)
    _tx(s, Inches(0.9), Inches(5.15), Inches(11.5), Inches(0.4),
        "Supervisor:   Mr. Varunendra Pratap Singh",
        size=Pt(15), bold=True, color=WHITE)
    _tx(s, Inches(0.9), Inches(5.6), Inches(11.5), Inches(0.4),
        "Principal Software Engineer  ·  Walmart Global Tech, Bengaluru",
        size=Pt(12), color=LIGHT_BLUE)
    _tx(s, Inches(0.9), Inches(6.15), Inches(11.5), Inches(0.4),
        "Candidate:   Vishal Singh   ·   2020AA05641",
        size=Pt(15), bold=True, color=WHITE)
    _tx(s, Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.4),
        "M.Tech (AI & ML)  ·  Birla Institute of Technology & Science, Pilani (WILP)",
        size=Pt(12), color=LIGHT_BLUE)
    _tx(s, Inches(0.9), Inches(7.05), Inches(11.5), Inches(0.3),
        "August 2026", size=Pt(10), color=MED_GRAY)

    # 2 · Agenda ─────────────────────────────────────────────────────────────
    s = new()
    _header_bar(s, "Agenda", "Focus of today — post mid-semester work only")
    items = [
        ("1", "Since Mid-Sem — What's New",         "Post-midsem workstream map"),
        ("2", "Review & Validate  (HITL)",           "Human-in-the-loop correction workflow"),
        ("3", "Smart Reply  —  Triple Draft",        "GPT-4o  +  Mistral 7B  +  Smart Composer"),
        ("4", "Smart Reply  —  Prompt & Few-Shot",   "Prompt template + past validated replies"),
        ("5", "Smart Reply  —  Worked Example",      "One complaint  →  three drafts"),
        ("6", "Learning Loop",                       "Corrections + posted replies → retraining"),
        ("7", "Post Explorer",                       "Multi-facet search across analysed posts"),
        ("8", "Post Lifecycle  (Kanban)",            "Triage → Acknowledged → In-Progress → Resolved"),
        ("9", "Insights & Competitor Analysis",      "Strategic weekly summaries + cross-brand"),
        ("10", "Notification Centre",                "In-app + email + Slack, group-routed"),
        ("11", "ModernBERT — Final Results",         "Macro-F1  0.6272 → 0.7642  (+13.7 pts)"),
        ("12", "Vision — Multi-Pass Payoff",         "Hallucination  50% → 0%,  extraction 25% → 75%"),
        ("13", "Trust-Score Evaluation",             "15% flagged; 12 of 15 confirmed by annotator"),
        ("14", "Live Demo",                          "Dashboard walkthrough (screenshots)"),
        ("15", "Conclusions & Future Work",          "RQ1–RQ4 outcomes + realistic roadmap"),
    ]
    x0, y0 = Inches(0.5), Inches(1.05)
    col_w, row_h = Inches(6.15), Inches(0.62)
    for i, (n, title, sub) in enumerate(items):
        col = i // 8
        row = i % 8
        lx = x0 + col * (col_w + Inches(0.15))
        ly = y0 + row * (row_h + Inches(0.12))
        _rect(s, lx, ly, col_w, row_h, LIGHT_BLUE, line=WALMART_BLUE)
        _tx(s, lx + Inches(0.1), ly, Inches(0.6), row_h,
            n, size=Pt(14), bold=True, color=WALMART_BLUE,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        _tx(s, lx + Inches(0.75), ly + Inches(0.07), col_w - Inches(0.85), Inches(0.28),
            title, size=Pt(12), bold=True, color=DARK_BLUE)
        _tx(s, lx + Inches(0.75), ly + Inches(0.34), col_w - Inches(0.85), Inches(0.26),
            sub, size=Pt(10), color=MED_GRAY)
    _footer(s, page[0], TOTAL)

    # 3 · What's new since mid-sem ──────────────────────────────────────────
    s = new()
    _header_bar(s, "Since Mid-Sem — What's New",
                "Mid-sem shipped the pipeline; post-midsem made it usable for analysts")
    _tx(s, Inches(0.5), Inches(1.05), Inches(6.0), Inches(0.35),
        "Delivered at mid-sem", size=Pt(13), bold=True, color=AMBER)
    _bullets(s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(3.8), [
        "6-layer offline-first pipeline (25 subreddits)",
        "ModernBERT fine-tuning up to Stage 2  (F1 0.7285)",
        "Multi-pass vision with 0% hallucination on pilot",
        "Trust-score design + weighted formula",
        "Brand-Health + Alert Feed dashboard pages",
        "Group-routed email / Slack notifications",
    ], size=Pt(12))
    _tx(s, Inches(6.85), Inches(1.05), Inches(6.0), Inches(0.35),
        "Added post-midsem  (this deck)", size=Pt(13), bold=True, color=GREEN)
    _bullets(s, Inches(6.85), Inches(1.4), Inches(6.0), Inches(3.8), [
        "Review & Validate  —  human-in-the-loop UI + backend",
        "Smart Reply Composer  —  dual-draft (rule + FLAN-T5)",
        "Learning-loop store  —  corrections & posted replies",
        "Post Explorer  —  multi-facet search",
        "Post Lifecycle  —  Kanban with 2-step Resolve",
        "Insights & Competitor Analysis page",
        "In-app Notification Centre mirroring the digest",
        "ModernBERT Stage 3 (final)  —  Macro-F1 0.7285 → 0.7642  (5-fold OOF)",
        "Trust-score end-to-end evaluation (n = 200)",
    ], size=Pt(12))
    _rect(s, Inches(0.5), Inches(5.55), Inches(12.3), Inches(1.35),
          LIGHT_BLUE, line=WALMART_BLUE)
    _tx(s, Inches(0.7), Inches(5.7), Inches(12.0), Inches(0.4),
        "Framing", size=Pt(13), bold=True, color=WALMART_BLUE)
    _tx(s, Inches(0.7), Inches(6.1), Inches(12.0), Inches(0.8),
        "Mid-sem answered can the pipeline produce trust-weighted sentiment?    "
        "Post-midsem answers can an analyst act on it, correct it, and feed the corrections back into the model?",
        size=Pt(12), color=DARK_GRAY)
    _footer(s, page[0], TOTAL)

    # 4 · Review & Validate ─────────────────────────────────────────────────
    s = new()
    _header_bar(s, "Review & Validate  —  Human-in-the-Loop",
                "Every prediction can be corrected; every correction is logged")
    _bullets(s, Inches(0.5), Inches(1.05), Inches(6.0), Inches(2.8), [
        "Queue sorted by  priority (P1 first)  ·  trust × confidence",
        "Analyst can override sentiment and aspect tags in one click",
        "Corrections written to  feedback table  with analyst + timestamp",
        "Generate Drafts  →  three reply options  (GPT-4o + Mistral + Smart Composer)",
        "Analyst picks A, B or C, edits inline, and posts to Reddit",
        "Posted replies also captured  →  future few-shot pool",
    ], size=Pt(12))
    _image(s, UI_DIR / "review_validate.png",
           Inches(6.75), Inches(1.05), width=Inches(6.2))
    _rect(s, Inches(0.5), Inches(5.7), Inches(6.0), Inches(1.2),
          GREEN_TINT, line=GREEN)
    _tx(s, Inches(0.7), Inches(5.85), Inches(5.6), Inches(0.35),
        "Why it matters", size=Pt(12), bold=True, color=GREEN)
    _tx(s, Inches(0.7), Inches(6.2), Inches(5.6), Inches(0.7),
        "Turns a passive dashboard into a training signal — corrections harvested here "
        "drive the ModernBERT re-training loop and reply few-shot pool.",
        size=Pt(11), color=DARK_GRAY)
    _footer(s, page[0], TOTAL)

    # 5 · Smart Reply Composer — Triple Draft ──────────────────────────────
    s = new()
    _header_bar(s, "Smart Reply Composer  —  Triple Draft",
                "One prompt, three generators; analyst picks the best draft")
    drafts = [
        ("Draft A", "GPT-4o", "Walmart LLM Gateway",
         "Strong reasoning, safest tone\nGuardrails via Walmart proxy\n~$0.0002 / reply\nFallback: direct OpenAI",
         WALMART_BLUE),
        ("Draft B", "Mistral 7B-Instruct", "Local Ollama  (:11434)",
         "Open-weights, free\nRuns fully offline\n~15 s warm latency\nStrong at retail jargon",
         PURPLE),
        ("Draft C", "Smart Composer", "Deterministic template",
         "No LLM, always available\nZero latency, zero cost\nContent-aware phrase pools\nSafety-net fallback",
         GREEN),
    ]
    x = Inches(0.4)
    y = Inches(1.05)
    w = Inches(4.15)
    h = Inches(3.35)
    for i, (label, model, transport, body, col) in enumerate(drafts):
        _rect(s, x, y, w, h, LIGHT_GRAY, line=col)
        _rect(s, x, y, w, Inches(0.5), col)
        _tx(s, x, y + Inches(0.08), w, Inches(0.35),
            label, size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _tx(s, x, y + Inches(0.6), w, Inches(0.4),
            model, size=Pt(16), bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
        _tx(s, x, y + Inches(1.05), w, Inches(0.35),
            transport, size=Pt(10), color=MED_GRAY, align=PP_ALIGN.CENTER)
        _tx(s, x + Inches(0.2), y + Inches(1.5), w - Inches(0.4), Inches(1.8),
            body, size=Pt(11), color=DARK_GRAY, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE)
        x = x + w + Inches(0.15)
    # Cascade + fallback logic
    _rect(s, Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.3),
          LIGHT_BLUE, line=WALMART_BLUE)
    _tx(s, Inches(0.7), Inches(4.75), Inches(12.0), Inches(0.35),
        "Fallback logic  —  no draft is ever missing", size=Pt(13), bold=True,
        color=WALMART_BLUE)
    _bullets(s, Inches(0.7), Inches(5.1), Inches(12.0), Inches(1.7), [
        "GPT-4o unavailable  (no gateway key / consumer-ID / network)  →  Smart Composer draft in slot A",
        "Mistral unavailable  (Ollama not running)  →  Smart Composer draft in slot B",
        "Smart Composer  always  produces slot C  —  guaranteed reply",
        "UI labels each card with the actual model used, plus an offline-fallback badge when applicable",
    ], size=Pt(11))
    _footer(s, page[0], TOTAL)

    # 6 · Smart Reply — Prompt Design & Few-Shot ────────────────────────────
    s = new()
    _header_bar(s, "Smart Reply  —  Prompt Design & Few-Shot",
                "Same prompt fed to GPT and Mistral; Smart Composer skips the LLM step")
    _tx(s, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.35),
        "Prompt template  (built in  _build_reply_prompt)",
        size=Pt(13), bold=True, color=DARK_BLUE)
    prompt_box_top = Inches(1.4)
    _rect(s, Inches(0.5), prompt_box_top, Inches(12.3), Inches(3.55),
          RGBColor(0x1E, 0x29, 0x3B), line=DARK_BLUE)
    prompt_text = (
        "You are a senior Walmart customer-care analyst replying on Reddit.\n"
        "Write ONE reply to the customer below. Keep it 2-4 sentences,\n"
        "empathetic, specific to their complaint, no corporate jargon,\n"
        "no hashtags, no emojis. Do NOT promise refunds you can't verify;\n"
        "invite them to DM order details if action is needed. Sign off as\n"
        "a real person, not a brand.\n"
        "\n"
        "Example customer post: {past_post_1}\n"
        "Example analyst reply: {past_reply_1}\n"
        "Example customer post: {past_post_2}\n"
        "Example analyst reply: {past_reply_2}\n"
        "\n"
        "Subreddit: r/{subreddit}\n"
        "Customer ({author}) complaint about: {aspect_1}, {aspect_2}\n"
        "Customer post:\n{post_title + post_body, ≤ 1200 chars}\n"
        "\n"
        "Reply:"
    )
    _tx(s, Inches(0.7), prompt_box_top + Inches(0.15), Inches(11.9), Inches(3.25),
        prompt_text, size=Pt(10), color=RGBColor(0xE6, 0xED, 0xF7))
    # Few-shot source
    _rect(s, Inches(0.5), Inches(5.15), Inches(6.0), Inches(1.75),
          GREEN_TINT, line=GREEN)
    _tx(s, Inches(0.7), Inches(5.25), Inches(5.7), Inches(0.35),
        "Few-shot source  —  feedback  table",
        size=Pt(12), bold=True, color=GREEN)
    _bullets(s, Inches(0.7), Inches(5.6), Inches(5.7), Inches(1.25), [
        "Top-3 past validated replies for the same aspect",
        "Written by human analysts, already posted to Reddit",
        "Refreshed on every  Generate  click  →  adapts over time",
    ], size=Pt(10))
    # Guardrails
    _rect(s, Inches(6.75), Inches(5.15), Inches(6.05), Inches(1.75),
          AMBER_TINT, line=AMBER)
    _tx(s, Inches(6.95), Inches(5.25), Inches(5.7), Inches(0.35),
        "Style guardrails baked into the prompt",
        size=Pt(12), bold=True, color=AMBER)
    _bullets(s, Inches(6.95), Inches(5.6), Inches(5.7), Inches(1.25), [
        "2 – 4 sentences, empathetic, specific",
        "No corporate jargon, no hashtags, no emojis",
        "No unverifiable refund promises  →  DM to continue",
        "Sign off as a person, not a brand",
    ], size=Pt(10))
    _footer(s, page[0], TOTAL)

    # 7 · Smart Reply — Worked Example ──────────────────────────────────────
    s = new()
    _header_bar(s, "Smart Reply  —  Worked Example  (all three drafts real, live-captured)",
                "raw_posts.id = reddit_1u2bgdw  ·  r/samsclub  ·  gateway ✓  ·  Ollama ✓  ·  Smart Composer ✓")
    # Customer complaint card
    _rect(s, Inches(0.5), Inches(1.0), Inches(12.3), Inches(2.05),
          AMBER_TINT, line=AMBER)
    _tx(s, Inches(0.7), Inches(1.08), Inches(12.0), Inches(0.35),
        "CUSTOMER COMPLAINT  ·  r/samsclub", size=Pt(11), bold=True, color=AMBER)
    _tx(s, Inches(0.7), Inches(1.4), Inches(12.0), Inches(0.32),
        "Why do you guys sell whole pizzas made hours ago to customers?",
        size=Pt(12), bold=True, color=DARK_BLUE)
    _tx(s, Inches(0.7), Inches(1.72), Inches(12.0), Inches(0.9),
        "\"The very few times I've gotten a whole pie, it'll be stuff premade and left in the hot case for like "
        "an hour 30mins before it's in my hand. How can I tell? They put a sticker with the date and the pizza "
        "looks and taste hours old. Pizza is meant to be made to order and waited for.\"",
        size=Pt(10), color=DARK_GRAY)
    _tx(s, Inches(0.7), Inches(2.6), Inches(8.5), Inches(0.28),
        "ModernBERT: negative (conf 0.9999997)  ·  aspects: customer service · product quality · store experience  ·  trust 0.66",
        size=Pt(9), color=MED_GRAY)
    _link(s, Inches(9.25), Inches(2.6), Inches(3.5), Inches(0.28),
          "https://www.reddit.com/r/samsclub/comments/1u2bgdw/why_do_you_guys_sell_whole_pizzas_made_hours_ago/",
          label="🔗 Open original on Reddit",
          size=Pt(9), bold=True, align=PP_ALIGN.RIGHT)
    # Three drafts side-by-side (VERBATIM from the live capture with all 3 engines running)
    drafts = [
        ("A · GPT-4o  (Walmart Gateway)", WALMART_BLUE, RGBColor(0xE8, 0xF4, 0xFD),
         "Hi u/there, I'm sorry to hear about your experience with the pizzas. I understand "
         "how disappointing it must be to get something that feels less fresh than expected. "
         "While the hot case is meant to provide convenience, I agree that quality should "
         "never be sacrificed. If you'd like, feel free to DM me with the details of your "
         "store location so we can look into this and address it properly. – [Your Name]"),
        ("B · Mistral 7B  (local Ollama)", PURPLE, PURPLE_TINT,
         "Hi u/there, I appreciate your concern about the quality of our prepared pizzas. "
         "It's important to us that our customers receive fresh and delicious food. The "
         "stickers with dates are there to help us manage the rotation of our products, "
         "ensuring they are consumed before they become stale. If you ever encounter a "
         "pizza that doesn't meet your expectations, please feel free to DM me the details "
         "of your order so I can look into it further. Thank you for taking the time to "
         "share your feedback."),
        ("C · Smart Composer  (no-LLM)", GREEN, GREEN_TINT,
         "Hi u/there — thanks for flagging this. completely understand the frustration "
         "around your order. Drop us a private message with the order # or store info and "
         "we'll take it from here. — Walmart Care"),
    ]
    x = Inches(0.4)
    y = Inches(3.05)
    w = Inches(4.15)
    h = Inches(3.15)
    for label, col, tint, body in drafts:
        _rect(s, x, y, w, h, tint, line=col)
        _rect(s, x, y, w, Inches(0.42), col)
        _tx(s, x, y + Inches(0.05), w, Inches(0.35),
            label, size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _tx(s, x + Inches(0.15), y + Inches(0.5), w - Inches(0.3), h - Inches(0.6),
            body, size=Pt(9), color=DARK_GRAY, anchor=MSO_ANCHOR.TOP)
        x = x + w + Inches(0.15)
    # Bottom explainer bar
    _rect(s, Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.6),
          LIGHT_BLUE, line=WALMART_BLUE)
    _tx(s, Inches(0.7), Inches(6.4), Inches(12.0), Inches(0.5),
        "One prompt fed to all three engines.  Analyst picks the best draft, edits inline, clicks Save & Open Reddit  →  "
        "the posted reply is written to  feedback  and becomes the next post's top few-shot example.",
        size=Pt(10), color=DARK_GRAY, anchor=MSO_ANCHOR.MIDDLE)
    _footer(s, page[0], TOTAL)

    # 6 · Learning Loop ─────────────────────────────────────────────────────
    s = new()
    _header_bar(s, "Learning Loop  —  Feedback → Retraining",
                "Every human action becomes a training signal  —  hot loop (live) + warm loop (monthly)")
    captures = [
        ("Label / aspect correction", "Review & Validate", WALMART_BLUE),
        ("Trust-score override",      "Review & Validate", PURPLE),
        ("Posted reply",              "Smart Reply Composer", GREEN),
        ("Lifecycle transition",      "Kanban board", AMBER),
    ]
    x = Inches(0.4)
    y = Inches(1.05)
    w = Inches(3.05)
    h = Inches(1.15)
    for title, src, col in captures:
        _rect(s, x, y, w, h, LIGHT_GRAY, line=col)
        _rect(s, x, y, w, Inches(0.35), col)
        _tx(s, x, y + Inches(0.05), w, Inches(0.3),
            title, size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _tx(s, x, y + Inches(0.5), w, Inches(0.65),
            src, size=Pt(11), color=DARK_GRAY, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE)
        x = x + w + Inches(0.15)
    _rect(s, Inches(0.5), Inches(2.5), Inches(12.3), Inches(0.7),
          DARK_BLUE, line=DARK_BLUE)
    _tx(s, Inches(0.7), Inches(2.55), Inches(12.0), Inches(0.6),
        "feedback  table  —  one row per human action, JSON payload, partition_key = analyst_id",
        size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE)
    _rect(s, Inches(0.5), Inches(3.4), Inches(6.05), Inches(2.5),
          GREEN_TINT, line=GREEN)
    _tx(s, Inches(0.7), Inches(3.5), Inches(5.7), Inches(0.35),
        "HOT LOOP  ·  in-context, live", size=Pt(13), bold=True, color=GREEN)
    _bullets(s, Inches(0.7), Inches(3.9), Inches(5.7), Inches(1.95), [
        "Every Generate Drafts click queries feedback",
        "SELECT … WHERE kind = 'auto_reply_posted'",
        "ORDER BY created_at DESC LIMIT 5  →  top 3 pairs",
        "Injected into the prompt as few-shot examples",
        "Effect visible on the next reply  —  no training",
    ], size=Pt(11))
    _rect(s, Inches(6.75), Inches(3.4), Inches(6.05), Inches(2.5),
          PURPLE_TINT, line=PURPLE)
    _tx(s, Inches(6.95), Inches(3.5), Inches(5.7), Inches(0.35),
        "WARM LOOP  ·  supervised, monthly", size=Pt(13), bold=True, color=PURPLE)
    _bullets(s, Inches(6.95), Inches(3.9), Inches(5.7), Inches(1.95), [
        "Export corrections since last train date",
        "Append to Walmart-200  →  Walmart-N (augmented)",
        "Rerun ModernBERT Stage-3 with 5-fold OOF CV",
        "New checkpoint wins only if OOF F1 improves",
        "Symlink flipped in config/models.yaml",
    ], size=Pt(11))
    _rect(s, Inches(0.5), Inches(6.05), Inches(12.3), Inches(0.85),
          LIGHT_BLUE, line=WALMART_BLUE)
    _tx(s, Inches(0.7), Inches(6.15), Inches(12.0), Inches(0.35),
        "Two failure modes eliminated by design",
        size=Pt(12), bold=True, color=WALMART_BLUE)
    _tx(s, Inches(0.7), Inches(6.5), Inches(12.0), Inches(0.35),
        "(1)  hot loop uses in-context learning  —  no retrain, no catastrophic forgetting.   "
        "(2)  warm loop uses 5-fold OOF CV  —  no leakage between train and eval sets.",
        size=Pt(11), color=DARK_GRAY)
    _footer(s, page[0], TOTAL)

    # 7 · Post Explorer ─────────────────────────────────────────────────────
    s = new()
    _header_bar(s, "Post Explorer  —  Multi-Facet Search",
                "Find the exact 20 posts an analyst needs to look at today")
    _tx(s, Inches(0.5), Inches(1.05), Inches(6.0), Inches(0.35),
        "Facets", size=Pt(13), bold=True, color=DARK_BLUE)
    _bullets(s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(3.8), [
        "Sentiment  ·  neg / neu / pos",
        "Confidence slider  (0.50 – 1.00)",
        "Trust-score slider  (0.00 – 1.00)",
        "Subreddit  —  multi-select from 25 tracked",
        "Aspect  —  8-item retail taxonomy",
        "Date range  —  today / week / month / custom",
        "Full-text search  —  title + body",
    ], size=Pt(12))
    _image(s, UI_DIR / "post_explorer.png",
           Inches(6.75), Inches(1.05), width=Inches(6.2))
    _rect(s, Inches(0.5), Inches(5.55), Inches(12.3), Inches(1.35),
          LIGHT_GRAY, line=WALMART_BLUE)
    _tx(s, Inches(0.7), Inches(5.7), Inches(12.0), Inches(0.35),
        "Per-post card", size=Pt(12), bold=True, color=DARK_BLUE)
    _tx(s, Inches(0.7), Inches(6.05), Inches(12.0), Inches(0.8),
        "Title + body excerpt · sentiment badge with confidence % · trust tier (H / M / L) · aspect chips · "
        "subreddit + timestamp · actions:  Review  ·  Add to Lifecycle  ·  View Details.",
        size=Pt(11), color=DARK_GRAY)
    _footer(s, page[0], TOTAL)

    # 8 · Post Lifecycle Kanban ─────────────────────────────────────────────
    s = new()
    _header_bar(s, "Post Lifecycle  —  Kanban Workflow",
                "Track a complaint from Triage to Resolved with SLA visibility")
    states = [
        ("TRIAGED",       "New P1 / P2\nposts land here",       AMBER),
        ("ACKNOWLEDGED",  "Assigned\nto analyst",               WALMART_BLUE),
        ("IN PROGRESS",   "Reply being\ndrafted / posted",      PURPLE),
        ("RESOLVED",      "Reply posted or\nno-action closed",  GREEN),
    ]
    x = Inches(0.5)
    y = Inches(1.35)
    w = Inches(3.0)
    h = Inches(1.9)
    for i, (t, body, col) in enumerate(states):
        _rect(s, x, y, w, h, LIGHT_GRAY, line=col)
        _rect(s, x, y, w, Inches(0.45), col)
        _tx(s, x, y + Inches(0.05), w, Inches(0.35),
            t, size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _tx(s, x, y + Inches(0.55), w, h - Inches(0.6),
            body, size=Pt(12), color=DARK_GRAY, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE)
        if i < 3:
            _arrow(s, x + w + Inches(0.005),
                   y + Inches(0.75), Inches(0.15), Inches(0.4))
        x = x + w + Inches(0.15)
    _image(s, UI_DIR / "lifecycle_kanban.png",
           Inches(0.5), Inches(3.55), width=Inches(7.5))
    _rect(s, Inches(8.3), Inches(3.55), Inches(4.5), Inches(3.15),
          LIGHT_BLUE, line=WALMART_BLUE)
    _tx(s, Inches(8.5), Inches(3.7), Inches(4.2), Inches(0.35),
        "Resolve modal  —  2-step flow", size=Pt(12), bold=True, color=WALMART_BLUE)
    _bullets(s, Inches(8.5), Inches(4.05), Inches(4.2), Inches(2.6), [
        "Step 1 · save action note + optional LLM reply",
        "(a) Save & open Reddit  →  reply copied to clipboard",
        "(b) OR  Resolve (no-reply)  →  close without posting",
        "Step 2 · paste on Reddit → Mark Resolved",
        "All transitions timestamped  →  SLA analytics",
    ], size=Pt(11))
    _footer(s, page[0], TOTAL)

    # 9 · Insights & Competitor ─────────────────────────────────────────────
    s = new()
    _header_bar(s, "Insights & Competitor Analysis",
                "Strategic view — issue rankings, cross-brand pulse, LLM summaries")
    features = [
        ("Priority-Negatives",     "Top issues ranked by\nvolume × severity × recency\nby aspect",       WALMART_BLUE),
        ("Competitor Pulse",       "Walmart vs Costco,\nTarget, Amazon on shared aspects",              PURPLE),
        ("Weekly LLM Summaries",   "Natural-language weekly\ndigest + action items\n+ emerging topics", GREEN),
        ("Aspect Drilldown",       "8-aspect taxonomy\nper-aspect sentiment trend\n+ representative posts", AMBER),
    ]
    positions = [(0.5, 1.05), (6.9, 1.05), (0.5, 3.55), (6.9, 3.55)]
    for (t, body, col), (lx, ly) in zip(features, positions):
        _rect(s, Inches(lx), Inches(ly), Inches(6.0), Inches(2.35),
              LIGHT_GRAY, line=col)
        _rect(s, Inches(lx), Inches(ly), Inches(6.0), Inches(0.45), col)
        _tx(s, Inches(lx), Inches(ly + 0.06), Inches(6.0), Inches(0.35),
            t, size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _tx(s, Inches(lx + 0.2), Inches(ly + 0.55), Inches(5.6), Inches(1.7),
            body, size=Pt(12), color=DARK_GRAY, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE)
    _tx(s, Inches(0.5), Inches(6.15), Inches(12.3), Inches(0.75),
        "Screenshot on Slide 17 (Live Demo grid) — insights_competitor page.",
        size=Pt(10), color=MED_GRAY, align=PP_ALIGN.CENTER)
    _footer(s, page[0], TOTAL)

    # 10 · Notification Centre ──────────────────────────────────────────────
    s = new()
    _header_bar(s, "Notification Centre",
                "In-app mirror of the group-routed email / Slack digest")
    _bullets(s, Inches(0.5), Inches(1.05), Inches(6.0), Inches(3.5), [
        "Every P1 / P2 email or Slack alert is mirrored in the app",
        "Notifications grouped by  subreddit ↔ business team",
        "Read / unread state persisted per analyst",
        "Deep-links open the post directly in Review & Validate",
        "Powered by the same  alerts  rows that drive email/Slack",
    ], size=Pt(12))
    _image(s, UI_DIR / "notifications.png",
           Inches(6.75), Inches(1.05), width=Inches(6.2))
    _rect(s, Inches(0.5), Inches(5.05), Inches(6.0), Inches(1.8),
          LIGHT_BLUE, line=WALMART_BLUE)
    _tx(s, Inches(0.7), Inches(5.2), Inches(5.6), Inches(0.35),
        "Priority thresholds", size=Pt(12), bold=True, color=WALMART_BLUE)
    _tx(s, Inches(0.7), Inches(5.55), Inches(5.6), Inches(1.25),
        "P1 — trust ≥ 0.70  and  confidence ≥ 0.80\nP2 — trust ≥ 0.50  and  confidence ≥ 0.60\nSource: config/models.yaml + dispatcher.py",
        size=Pt(11), color=DARK_GRAY)
    _footer(s, page[0], TOTAL)

    # 11 · ModernBERT Final Results ─────────────────────────────────────────
    s = new()
    _header_bar(s, "ModernBERT  —  Final Training Results",
                "3-stage curriculum,  5-fold out-of-fold CV,  n = 200")
    _kpi_tile(s, Inches(0.5), Inches(1.05), Inches(4.0), Inches(1.75),
              "Overall Macro-F1", "0.6272 → 0.7642",
              value_color=GREEN, tint=GREEN_TINT, value_size=Pt(24))
    _kpi_tile(s, Inches(4.65), Inches(1.05), Inches(4.0), Inches(1.75),
              "Long posts  (≥ 512 tok,  n = 7)", "5 / 7  →  7 / 7  correct",
              value_color=GREEN, tint=GREEN_TINT, value_size=Pt(20))
    _kpi_tile(s, Inches(8.8), Inches(1.05), Inches(4.0), Inches(1.75),
              "Uplift over baseline", "+13.7 pts",
              value_color=WALMART_BLUE, value_size=Pt(28))
    _tx(s, Inches(0.5), Inches(3.0), Inches(12.3), Inches(0.35),
        "3-stage curriculum", size=Pt(13), bold=True, color=DARK_BLUE)
    _table(s, Inches(0.5), Inches(3.4), Inches(12.3), Inches(1.85), [
        ["Stage", "Data",                                          "Macro-F1", "Δ vs baseline"],
        ["Baseline (Twitter-RoBERTa)",  "off-the-shelf",           "0.6272",   "—"],
        ["Stage 1",   "TweetEval sentiment (~45 k tweets)",        "0.6810",   "+5.4"],
        ["Stage 2",   "GoEmotions-3class Reddit (~54 k, pseudo)",  "0.7285",   "+10.1"],
        ["Stage 3  (final)", "Walmart-200 hand-labelled",          "0.7642",   "+13.7"],
    ], first_col_bold=True, highlight_rows={4: GREEN_TINT})
    _tx(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.35),
        "Per-length-bucket accuracy  (5-fold OOF predictions)",
        size=Pt(13), bold=True, color=DARK_BLUE)
    _table(s, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.1), [
        ["Bucket  (token count)",           "Baseline correct", "ModernBERT correct",  "Recovered"],
        ["Short-to-medium  (< 512, n = 193)", "138 / 193 (72 %)",  "159 / 193 (82 %)",     "+21"],
        ["Long  (≥ 512, n = 7, all negative)", "5 / 7",            "7 / 7",                "+2"],
    ], first_col_bold=True, highlight_rows={2: GREEN_TINT})
    _footer(s, page[0], TOTAL)

    # 12 · Vision Multi-Pass Payoff ─────────────────────────────────────────
    s = new()
    _header_bar(s, "Vision Pipeline  —  Multi-Pass Payoff",
                "Same model (Gemma 3 4B), same hardware, smarter calling strategy")
    passes = [
        ("PASS 1", "STRUCTURE",  "What TYPE of image?\nscreenshot / photo /\nreceipt / app / meme", PURPLE),
        ("PASS 2", "TILE",       "Split into 2–4 crops\n→ 2–4× effective\nresolution, no resize",   WALMART_BLUE),
        ("PASS 3", "EXTRACT",    "Per tile: read ALL\ntext verbatim\n(quotes, prices, UI)",         AMBER),
        ("PASS 4", "MERGE",      "Text-only LLM call\n— image REMOVED —\ncannot invent visuals",     GREEN),
    ]
    x = Inches(0.4)
    y = Inches(1.15)
    w = Inches(3.0)
    h = Inches(2.15)
    for i, (label, title, body, col) in enumerate(passes):
        _stage_card(s, x, y, w, h, label, title, body, col)
        if i < 3:
            _arrow(s, x + w + Inches(0.005),
                   y + Inches(0.8), Inches(0.18), Inches(0.4))
        x = x + w + Inches(0.15)
    _tx(s, Inches(0.5), Inches(3.65), Inches(12.3), Inches(0.35),
        "Evaluation on 32 retail screenshots", size=Pt(13), bold=True, color=DARK_BLUE)
    _table(s, Inches(0.5), Inches(4.05), Inches(12.3), Inches(2.6), [
        ["Metric",                    "Single-pass",  "Multi-pass",  "Change"],
        ["Hallucination rate",         "50 %",         "0 %",         "eliminated"],
        ["Text-extraction success",    "25 %",         "75 %",        "3× better"],
        ["Retail-signal recall",       "40 %",         "81 %",        "2× better"],
        ["Fabricated claims",          "many",         "0",           "eliminated"],
        ["Latency / image (warm)",     "~5 s",         "~15 s",       "3× (accepted)"],
    ], first_col_bold=True,
       highlight_rows={1: GREEN_TINT, 2: GREEN_TINT, 3: GREEN_TINT})
    _footer(s, page[0], TOTAL)

    # 13 · Trust-Score Evaluation ───────────────────────────────────────────
    s = new()
    _header_bar(s, "Trust-Score  —  End-to-End Evaluation",
                "On the 200-post gold set, does the trust filter catch bad posts?")
    _kpi_tile(s, Inches(0.5), Inches(1.05), Inches(4.0), Inches(1.9),
              "Low-trust share  (T < 0.4)", "15 %",
              value_color=AMBER, tint=AMBER_TINT, value_size=Pt(38))
    _kpi_tile(s, Inches(4.65), Inches(1.05), Inches(4.0), Inches(1.9),
              "Human-agreed  (annotator cross-check)", "12 / 15",
              value_color=GREEN, tint=GREEN_TINT, value_size=Pt(38))
    _kpi_tile(s, Inches(8.8), Inches(1.05), Inches(4.0), Inches(1.9),
              "Agreement rate", "80 %",
              value_color=WALMART_BLUE, value_size=Pt(38))
    _tx(s, Inches(0.5), Inches(3.15), Inches(12.3), Inches(0.35),
        "Weighted formula  (decomposition shown inline in dashboard)",
        size=Pt(13), bold=True, color=DARK_BLUE)
    _rect(s, Inches(0.5), Inches(3.55), Inches(12.3), Inches(0.85),
          LIGHT_BLUE, line=WALMART_BLUE)
    _tx(s, Inches(0.7), Inches(3.7), Inches(12.0), Inches(0.55),
        "trust_score  =  0.4 × metadata  +  0.3 × dedup  +  0.3 × llm_credibility",
        size=Pt(18), bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE)
    _rect(s, Inches(0.5), Inches(4.65), Inches(12.3), Inches(2.15),
          GREEN_TINT, line=GREEN)
    _tx(s, Inches(0.7), Inches(4.8), Inches(12.0), Inches(0.35),
        "Design principle — flag, don't drop", size=Pt(13), bold=True, color=GREEN)
    _tx(s, Inches(0.7), Inches(5.15), Inches(12.0), Inches(1.6),
        "Low-trust posts are surfaced with an explanatory chip (\"promotional / duplicate / thin-metadata\") and can be "
        "overridden by an analyst in one click — the override is captured for the learning loop. "
        "Every constant in the formula has a stakeholder-arguable English rationale in config/models.yaml.",
        size=Pt(12), color=DARK_GRAY)
    _footer(s, page[0], TOTAL)

    # 14 · Evaluation summary ───────────────────────────────────────────────
    s = new()
    _header_bar(s, "Evaluation Summary  —  Post-Midsem Numbers",
                "All results reported against the same 200-post retail-Reddit gold set")
    _table(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(5.6), [
        ["Area",             "Metric",                                          "Value",         "Notes"],
        ["Sentiment",        "Macro-F1 (final)",                                "0.7642",        "+13.7 pts vs RoBERTa baseline"],
        ["Sentiment",        "Long-bucket accuracy  (≥ 512 tokens, n=7)",       "7 / 7 correct",  "Baseline 5 / 7  ·  all negative-class"],
        ["Sentiment",        "Cross-validation",                                "5-fold OOF",    "No test-set leakage"],
        ["Vision",           "Hallucination rate",                              "0 %",           "From 50 % on single-pass"],
        ["Vision",           "Text extraction success",                         "75 %",          "From 25 % on single-pass"],
        ["Vision",           "Retail-signal recall",                            "81 %",          "From 40 % on single-pass"],
        ["Trust",            "Low-trust share  (T < 0.4)",                      "15 %",          "Of the 200-post set"],
        ["Trust",            "Annotator agreement on low-trust",                "12 / 15",       "80 %"],
        ["Ops",              "Subreddits tracked",                              "25",            "Six community groups"],
        ["Ops",              "Scheduler cadence",                               "6 h + manual",  "asyncio lifespan task"],
        ["Ops",              "Dashboard pages shipped",                         "8",             "Brand Health · Alerts · Explorer · Review · Lifecycle · Insights · Pipeline · Notifications"],
    ], first_col_bold=True, body_font=Pt(10),
       highlight_rows={1: GREEN_TINT, 2: GREEN_TINT, 4: GREEN_TINT, 5: GREEN_TINT})
    _footer(s, page[0], TOTAL)

    # 15 · Contributions ────────────────────────────────────────────────────
    s = new()
    _header_bar(s, "Principal Contributions",
                "What this dissertation delivered — end to end")
    contribs = [
        ("C1", "Offline-first RSI pipeline",
         "Ingestion → trust → sentiment → aspects → vision → aggregation → alerts → dashboard.  Zero API cost, modular, deployable to Azure."),
        ("C2", "Fine-tuned ModernBERT (3-stage curriculum)",
         "Macro-F1 0.6272 → 0.7642 overall; 0.28 → 1.00 on long posts.  5-fold OOF cross-validation, offline training."),
        ("C3", "Multi-pass vision technique on Gemma 3 4B",
         "Hallucination 50 % → 0 %,  extraction 25 % → 75 %.  Adapts techniques from 5 recent VLM papers without vendor-blocked models."),
        ("C4", "Interpretable trust score + admission gate",
         "trust × confidence gate flags rather than drops low-credibility posts.  Every constant traceable to English rationale in config."),
        ("C5", "HITL learning-loop dashboard",
         "Review & Validate, Lifecycle Kanban, Insights, Notification centre.  Corrections & posted replies feed few-shot reply generation and future retraining."),
    ]
    y = Inches(1.05)
    for tag, title, body in contribs:
        _rect(s, Inches(0.5), y, Inches(12.3), Inches(1.05),
              LIGHT_GRAY, line=WALMART_BLUE)
        _rect(s, Inches(0.5), y, Inches(0.9), Inches(1.05), WALMART_BLUE)
        _tx(s, Inches(0.5), y, Inches(0.9), Inches(1.05),
            tag, size=Pt(20), bold=True, color=WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _tx(s, Inches(1.55), y + Inches(0.05), Inches(11.0), Inches(0.35),
            title, size=Pt(13), bold=True, color=DARK_BLUE)
        _tx(s, Inches(1.55), y + Inches(0.4), Inches(11.0), Inches(0.65),
            body, size=Pt(11), color=DARK_GRAY)
        y = y + Inches(1.14)
    _footer(s, page[0], TOTAL)

    # 16 · Live Demo (screenshot grid) ──────────────────────────────────────
    s = new()
    _header_bar(s, "Live Demo  —  Dashboard Walkthrough",
                "One post → alert → review → reply → lifecycle → resolved")
    shots = [
        ("Brand Health",       UI_DIR / "brand_health.png"),
        ("Alert Feed",         UI_DIR / "alert_feed.png"),
        ("Post Explorer",      UI_DIR / "post_explorer.png"),
        ("Review & Validate",  UI_DIR / "review_validate.png"),
        ("Lifecycle Kanban",   UI_DIR / "lifecycle_kanban.png"),
        ("Insights",           UI_DIR / "insights_competitor.png"),
    ]
    # 3 across × 2 down
    thumb_w = Inches(4.15)
    thumb_h = Inches(2.55)
    x0, y0 = Inches(0.35), Inches(1.05)
    gap = Inches(0.1)
    for i, (label, path) in enumerate(shots):
        col = i % 3
        row = i // 3
        lx = x0 + col * (thumb_w + gap)
        ly = y0 + row * (thumb_h + Inches(0.3))
        _image(s, path, lx, ly, width=thumb_w, height=thumb_h)
        _tx(s, lx, ly + thumb_h + Inches(0.02), thumb_w, Inches(0.25),
            label, size=Pt(11), bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    _footer(s, page[0], TOTAL)

    # 17 · What worked / what didn't ────────────────────────────────────────
    s = new()
    _header_bar(s, "Conclusions  —  What Worked, What's Still Open",
                "Honest read on 5 months of post-midsem work")
    _rect(s, Inches(0.5), Inches(1.05), Inches(6.05), Inches(5.75),
          GREEN_TINT, line=GREEN)
    _tx(s, Inches(0.7), Inches(1.2), Inches(5.7), Inches(0.4),
        "What worked", size=Pt(15), bold=True, color=GREEN)
    _bullets(s, Inches(0.7), Inches(1.65), Inches(5.7), Inches(5.0), [
        "3-stage ModernBERT curriculum landed the sentiment win we set out for  (Macro-F1 0.62 → 0.76 on 5-fold OOF).",
        "Long-post recovery came for free once we switched off truncation  (≥1024-tok context):  5/7 → 7/7 correct on the ≥ 512-token bucket.",
        "Removing the image on the final vision merge step ended the 50 % hallucination problem  —  the mechanism, not just the number, is what makes it credible.",
        "Trust score as flag-not-drop  matched the human annotator on 12 / 15 low-trust posts. Analysts trust it because they can override it.",
        "HITL feedback table quietly turned into the most useful piece of infrastructure  —  drives few-shot on every click today, retraining tomorrow.",
    ], size=Pt(11))
    _rect(s, Inches(6.75), Inches(1.05), Inches(6.05), Inches(5.75),
          AMBER_TINT, line=AMBER)
    _tx(s, Inches(6.95), Inches(1.2), Inches(5.7), Inches(0.4),
        "What's still open (honest)", size=Pt(15), bold=True, color=AMBER)
    _bullets(s, Inches(6.95), Inches(1.65), Inches(5.7), Inches(5.0), [
        "200-post gold set is small.  Numbers are OOF-CV, but a bigger blind held-out set is needed before I'd claim generalisation.",
        "Vision eval is 32 images  —  enough to sanity-check the mechanism, not enough to claim production-grade quality.",
        "Long bucket has only 7 posts and they are all negative-class  —  7/7 correct is evidence the truncation ceiling is gone, not proof of long-text mastery.",
        "Reply drafts still need an analyst edit  —  we track edit-distance but haven't shown it dropping consistently yet.",
        "Everything runs on one Mac. Multi-analyst concurrency and retraining automation aren't done  —  they're in the roadmap on the next slide.",
    ], size=Pt(11))
    _footer(s, page[0], TOTAL)

    # 18 · Future work (grounded, 3-horizon) ────────────────────────────────
    s = new()
    _header_bar(s, "Future Work  —  Grounded 3-Horizon Roadmap",
                "Only items with clear next steps; not a wishlist")
    horizons = [
        ("Now → 1 month",
         "Grow the gold set to ~500 posts using the HITL feedback the team is already producing; rerun 5-fold OOF and check the F1 delta.",
         GREEN),
        ("1 → 3 months",
         "Wire the monthly ModernBERT retrain into a Cron / Airflow job so it stops being a manual notebook run; add promotion gate on OOF F1.",
         WALMART_BLUE),
        ("3 → 6 months",
         "Move the feedback table to a shared managed database so multiple analysts can work concurrently; extend ingestion beyond Reddit only if the demand from the analyst team is real.",
         PURPLE),
    ]
    y = Inches(1.15)
    for title, body, col in horizons:
        _rect(s, Inches(0.5), y, Inches(12.3), Inches(1.6),
              LIGHT_GRAY, line=col)
        _rect(s, Inches(0.5), y, Inches(0.25), Inches(1.6), col)
        _tx(s, Inches(0.9), y + Inches(0.15), Inches(11.4), Inches(0.4),
            title, size=Pt(15), bold=True, color=DARK_BLUE)
        _tx(s, Inches(0.9), y + Inches(0.65), Inches(11.4), Inches(0.9),
            body, size=Pt(12), color=DARK_GRAY)
        y = y + Inches(1.75)
    _rect(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.5),
          LIGHT_BLUE, line=WALMART_BLUE)
    _tx(s, Inches(0.7), Inches(6.45), Inches(12.0), Inches(0.4),
        "Deliberately excluded  —  ideas without a clear next step  (auto-reply gate, seasonal P1 forecast, bilingual taxonomy)  —  covered in the report's Future Work chapter.",
        size=Pt(10), color=DARK_GRAY, anchor=MSO_ANCHOR.MIDDLE)
    _footer(s, page[0], TOTAL)

    # 19 · Source Code + Deliverables ───────────────────────────────────────
    s = new()
    _header_bar(s, "Source Code & Deliverables")
    _rect(s, Inches(0.5), Inches(1.05), Inches(12.3), Inches(1.1),
          LIGHT_BLUE, line=WALMART_BLUE)
    _tx(s, Inches(0.7), Inches(1.15), Inches(12.0), Inches(0.35),
        "Repository", size=Pt(12), bold=True, color=WALMART_BLUE)
    _tx(s, Inches(0.7), Inches(1.55), Inches(12.0), Inches(0.55),
        "https://gecgithub01.walmart.com/v0s01jh/Retail_Sentiment_Intelligence",
        size=Pt(16), bold=True, color=DARK_BLUE, anchor=MSO_ANCHOR.MIDDLE)
    _table(s, Inches(0.5), Inches(2.35), Inches(12.3), Inches(4.5), [
        ["Deliverable",              "Path in repo",                                                   "Contents"],
        ["Final Report (PDF)",       "docs/Sem_4/final/FINAL_REPORT_VishalSingh_2020AA05641.pdf",       "10 chapters + appendices + refs"],
        ["Final Report (LaTeX)",     "docs/Sem_4/final/latex/",                                        "Reproducible xelatex source"],
        ["Pipeline core",             "src/pipeline.py, src/ingestion/, src/analysis/",                 "6-layer async pipeline"],
        ["Sentiment training",       "scripts/train_modernbert_sentiment.py",                          "3-stage curriculum runner"],
        ["Evaluation notebook",      "evaluation/trust_score_walmart200.ipynb",                        "Reproducible 5-fold OOF eval"],
        ["Dashboard",                "frontend/  (React + Vite + Tailwind)",                           "8 pages, live via WebSocket"],
        ["Reproduction",             "Appendix B of the report  ·  start.sh",                          "Clone → conda → start"],
    ], first_col_bold=True, body_font=Pt(10))
    _footer(s, page[0], TOTAL)

    # 20 · Thank You / Q&A ──────────────────────────────────────────────────
    s = new()
    _set_bg(s, DARK_BLUE)
    _tx(s, Inches(0.5), Inches(1.9), Inches(12.3), Inches(1.2),
        "Thank You", size=Pt(72), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _tx(s, Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.5),
        "Retail Sentiment Intelligence  ·  Post Mid-Semester Presentation",
        size=Pt(18), color=YELLOW, align=PP_ALIGN.CENTER)
    _tx(s, Inches(0.5), Inches(4.1), Inches(12.3), Inches(0.5),
        "Questions ?", size=Pt(28), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _rect(s, Inches(2.0), Inches(5.05), Inches(9.3), Inches(1.4),
          LIGHT_BLUE, line=YELLOW)
    _tx(s, Inches(2.2), Inches(5.2), Inches(8.9), Inches(0.35),
        "Repository", size=Pt(11), bold=True, color=WALMART_BLUE,
        align=PP_ALIGN.CENTER)
    _tx(s, Inches(2.2), Inches(5.55), Inches(8.9), Inches(0.5),
        "gecgithub01.walmart.com/v0s01jh/Retail_Sentiment_Intelligence",
        size=Pt(15), bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    _tx(s, Inches(2.2), Inches(6.05), Inches(8.9), Inches(0.35),
        "Report:  docs/Sem_4/final/FINAL_REPORT_VishalSingh_2020AA05641.pdf",
        size=Pt(11), color=MED_GRAY, align=PP_ALIGN.CENTER)
    _tx(s, Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.4),
        "Vishal Singh  ·  2020AA05641  ·  BITS Pilani (WILP)  ·  Walmart Global Tech, Bengaluru",
        size=Pt(12), color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"[OK] wrote {OUT}  ({page[0]} slides)")


if __name__ == "__main__":
    build()
